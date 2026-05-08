"""
Build a text-led Alagappa HMS board deck.

This variant intentionally avoids architecture diagrams and dramatic visuals.
It keeps the board story focused on cost savings, internal capability building,
resources needed, and a controlled approval decision.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


FOREST = RGBColor(0x1F, 0x43, 0x32)
INK = RGBColor(0x12, 0x17, 0x14)
CANVAS = RGBColor(0xFF, 0xFF, 0xFF)
FOG = RGBColor(0xF7, 0xF8, 0xF6)
RULE = RGBColor(0xE3, 0xE8, 0xE4)
COPPER = RGBColor(0xA7, 0x7A, 0x35)
TINT = RGBColor(0xE7, 0xF0, 0xEC)
DIMMED = RGBColor(0x5E, 0x66, 0x62)

DISPLAY_FONT = "Aptos Display"
UI_FONT = "Aptos"
MONO_FONT = "Aptos Mono"

ROOT = Path("/Users/apple/Projects/MedBrains/medbrains")
OUTPUT = ROOT / "docs" / "Alagappa_HMS_Board_text_points_cost_muscle_resources.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL_SLIDES = 8


def add_bg(slide, color=CANVAS):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font=UI_FONT,
    size=16,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Emu(0)
    frame.margin_top = frame.margin_bottom = Emu(0)
    frame.vertical_anchor = anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_runs(slide, x, y, w, h, runs, *, size=24):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Emu(0)
    frame.margin_top = frame.margin_bottom = Emu(0)
    paragraph = frame.paragraphs[0]
    for text, style in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = style.get("font", DISPLAY_FONT)
        run.font.size = Pt(style.get("size", size))
        run.font.color.rgb = style.get("color", INK)
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
    return box


def add_card(slide, x, y, w, h, fill=FOG, border=RULE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.03
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_eyebrow(slide, text):
    add_text(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(8.0),
        Inches(0.25),
        text.upper(),
        font=MONO_FONT,
        size=9,
        color=FOREST,
        bold=True,
    )


def add_footer(slide, num):
    add_text(
        slide,
        Inches(0.6),
        Inches(7.18),
        Inches(7.8),
        Inches(0.25),
        "Alagappa Hospital - In-house HMS board decision - May 2026",
        font=MONO_FONT,
        size=8,
        color=DIMMED,
    )
    add_text(
        slide,
        Inches(11.4),
        Inches(7.18),
        Inches(1.3),
        Inches(0.25),
        f"{num} / {TOTAL_SLIDES}",
        font=MONO_FONT,
        size=8,
        color=DIMMED,
        align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, x, y, w, items, *, size=14, color=INK, gap=0.42):
    for idx, item in enumerate(items):
        yy = y + Inches(gap * idx)
        add_text(slide, x, yy, Inches(0.18), Inches(0.28), "-", size=size, color=FOREST, bold=True)
        add_text(slide, x + Inches(0.25), yy, w - Inches(0.25), Inches(0.34), item, size=size, color=color)


def title_slide(slide, eyebrow, title, subtitle):
    add_bg(slide)
    add_eyebrow(slide, eyebrow)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.95), title, font=DISPLAY_FONT, size=35, color=INK, bold=True)
    add_text(slide, Inches(0.6), Inches(1.85), Inches(11.8), Inches(0.55), subtitle, size=15, color=DIMMED, italic=True)


def add_table(slide, x, y, w, h, rows, *, widths=None, font_size=12):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = shape.table
    if widths:
        for idx, width in enumerate(widths):
            table.columns[idx].width = width
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FOREST if r == 0 else (CANVAS if r % 2 else FOG)
            cell.margin_left = Emu(120000)
            cell.margin_right = Emu(120000)
            cell.margin_top = Emu(70000)
            cell.margin_bottom = Emu(70000)
            frame = cell.text_frame
            frame.clear()
            frame.word_wrap = True
            paragraph = frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = value
            run.font.name = UI_FONT
            run.font.size = Pt(11 if r == 0 else font_size)
            run.font.bold = r == 0 or c == 0
            run.font.color.rgb = CANVAS if r == 0 else INK
    return table


# Slide 1
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Decision frame",
    "Approve the in-house HMS path for cost control and capability building",
    "This is not a software purchase decision. It is a hospital operating-system decision.",
)
add_card(s, Inches(0.6), Inches(2.75), Inches(12.1), Inches(3.35), TINT)
add_text(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(0.45), "The board story in one page", font=DISPLAY_FONT, size=20, color=FOREST, bold=True)
add_bullets(
    s,
    Inches(0.95),
    Inches(3.65),
    Inches(11.1),
    [
        "Save money by avoiding large upfront licence, AMC, change-request and duplicate-work costs.",
        "Build internal technology muscle so Alagappa can adapt workflows without waiting on vendor roadmaps.",
        "Start safely with OPD + pharmacy pilot, daily defect review and clear go / no-go gates.",
        "Ask for a small focused team, basic infrastructure, doctor/operations champions and board governance.",
    ],
    size=14,
    gap=0.52,
)
add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.35), "Board line: buy if we must, but build if we want control.", size=13, color=COPPER, bold=True)
add_footer(s, 1)

# Slide 2
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Main point 1 - cost savings",
    "The saving is not only licence price. It is recurring operating friction.",
    "Vendor quotes show the visible cost. Hospital operations carry the hidden cost.",
)
columns = [
    (
        "Avoid vendor lock-in cost",
        [
            "No large upfront licence as the default path.",
            "No AMC escalation for every branch and module.",
            "No quotation cycle for every workflow change.",
            "No exit cost when the hospital wants its own data and rules.",
        ],
    ),
    (
        "Reduce daily operating cost",
        [
            "Less printing, registers, photocopying and MRD file movement.",
            "Less duplicate entry across OPD, billing, pharmacy, lab and records.",
            "Fewer fixed counters through tablets, queue displays and shared workstations.",
            "Offline camps keep work moving instead of paper re-entry later.",
        ],
    ),
    (
        "Keep future cost under control",
        [
            "New departments and branches reuse the same platform.",
            "Compliance changes are in our backlog, not vendor backlog.",
            "Integrations like ABDM, NHCX, DICOM and devices stay owned by Alagappa.",
            "Resale or sister-hospital deployment becomes possible later.",
        ],
    ),
]
for idx, (heading, bullets) in enumerate(columns):
    x = Inches(0.6 + idx * 4.1)
    add_card(s, x, Inches(2.7), Inches(3.85), Inches(3.9))
    add_text(s, x + Inches(0.25), Inches(2.95), Inches(3.35), Inches(0.5), heading, font=DISPLAY_FONT, size=17, color=FOREST, bold=True)
    add_bullets(s, x + Inches(0.25), Inches(3.65), Inches(3.3), bullets, size=11, gap=0.48)
add_footer(s, 2)

# Slide 3
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Main point 2 - build hospital technology muscle",
    "The software is useful. The internal capability is the real long-term asset.",
    "A hospital that can change its own HMS can move faster than a hospital waiting for tickets.",
)
muscle_rows = [
    ("Capability we build", "Why it matters to the hospital"),
    ("Workflow ownership", "Doctors, nurses, pharmacy and billing get changes from an internal team that understands Alagappa."),
    ("Regulatory response", "NABH, ABDM, NHCX, DPDPA and pharmacy safety changes can be handled without vendor delay."),
    ("Integration confidence", "Lab machines, payment, insurance, WhatsApp, DICOM/PACS and future devices become repeatable work."),
    ("Data advantage", "Management dashboards, audit evidence, clinical quality indicators and branch comparison become possible."),
    ("Future readiness", "AI assistants, clinical decision support, analytics and sister-hospital rollout become internal muscle, not outsourced magic."),
]
add_table(s, Inches(0.6), Inches(2.65), Inches(12.1), Inches(3.85), muscle_rows, widths=[Inches(3.2), Inches(8.9)], font_size=12)
add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3), "Board line: this is how Alagappa stops being only a buyer and becomes a builder.", size=12, color=COPPER, bold=True)
add_footer(s, 3)

# Slide 4
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Main point 3 - resources needed",
    "The ask is focused: people, infrastructure, governance and pilot time.",
    "We do not need a big IT department before proving the pilot.",
)
resource_rows = [
    ("Resource", "What we need now", "Why"),
    ("People", "Lead architect + 2 senior engineers + existing hardware support", "Build, secure, deploy and maintain the core system."),
    ("Hospital champions", "One OPD doctor, one pharmacist, one billing/MRD owner, one admin owner", "Validate real workflow daily and stop assumptions early."),
    ("Infrastructure", "Cloud server, database backups, VPN/admin access, local backup storage", "Run production safely with recovery and controlled access."),
    ("Devices for pilot", "A few tablets/shared PCs, queue display, barcode/label printer where needed", "Prove paperless OPD + pharmacy without buying hardware for all departments."),
    ("Governance", "Daily pilot review + weekly board/management checkpoint", "Keep scope controlled and expand only after evidence."),
    ("Budget envelope", "Salary + infrastructure + pilot devices + contingency", "Spend in stages; avoid large vendor lock-in before proof."),
]
add_table(s, Inches(0.6), Inches(2.35), Inches(12.1), Inches(4.6), resource_rows, widths=[Inches(2.1), Inches(5.0), Inches(5.0)], font_size=10)
add_footer(s, 4)

# Slide 5
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Make the decision safe",
    "Pilot first. Expand only after proof.",
    "This reduces the risk of building and also improves our vendor negotiation fallback.",
)
gates = [
    ("June", "Team + infrastructure ready", "deployment, backups, access control, data migration plan"),
    ("July", "OPD + pharmacy pilot", "parallel run, live defects, doctor/pharmacy feedback"),
    ("August", "Stability review", "add IPD/lab/radiology only after OPD + pharmacy are stable"),
    ("September", "Go-live decision", "go live, extend pilot, or use vendor fallback with stronger negotiation"),
]
for idx, (month, title, body) in enumerate(gates):
    y = Inches(2.45 + idx * 0.95)
    add_card(s, Inches(0.6), y, Inches(12.1), Inches(0.72))
    add_text(s, Inches(0.9), y + Inches(0.16), Inches(1.25), Inches(0.35), month, font=MONO_FONT, size=13, color=FOREST, bold=True)
    add_text(s, Inches(2.35), y + Inches(0.12), Inches(3.6), Inches(0.35), title, font=DISPLAY_FONT, size=15, color=FOREST, bold=True)
    add_text(s, Inches(6.05), y + Inches(0.16), Inches(6.2), Inches(0.35), body, size=12, color=INK)
add_text(s, Inches(0.85), Inches(6.55), Inches(11.5), Inches(0.35), "The board is approving a controlled path, not blind faith.", size=13, color=COPPER, bold=True)
add_footer(s, 5)

# Slide 6
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "What we should not make the slide about",
    "Avoid architecture diagrams in the main board discussion.",
    "Keep technical proof in backup. The board needs decision clarity.",
)
avoid_rows = [
    ("Do not lead with", "Lead with instead"),
    ("CRDT / WASM / Rust internals", "Offline camp continuity, data safety and lower cloud cost."),
    ("Architecture boxes and arrows", "Who owns the workflow and who fixes it when doctors need changes."),
    ("Long module lists", "OPD, pharmacy, billing, lab, IPD path to go-live."),
    ("Vendor criticism only", "Cost control, independence and proof-based rollout."),
    ("Future AI promise", "Build the internal data and engineering base that makes AI possible later."),
]
add_table(s, Inches(0.6), Inches(2.45), Inches(12.1), Inches(3.75), avoid_rows, widths=[Inches(5.8), Inches(6.3)], font_size=12)
add_text(s, Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.35), "Technical material stays available for Q&A, not as the main persuasion.", size=12, color=COPPER, bold=True)
add_footer(s, 6)

# Slide 7
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Board approvals needed",
    "Three approvals are enough to move.",
    "Each approval has a practical control attached.",
)
approvals = [
    ("1", "Approve direction", "In-house HMS as primary path. Vendor remains fallback and benchmark."),
    ("2", "Approve resources", "Two senior engineers, infrastructure, pilot devices and hospital champions."),
    ("3", "Approve governance", "OPD + pharmacy first, daily pilot review, weekly management checkpoint."),
]
for idx, (num, title, body) in enumerate(approvals):
    y = Inches(2.45 + idx * 1.15)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), y, Inches(0.62), Inches(0.62))
    circle.fill.solid()
    circle.fill.fore_color.rgb = FOREST
    circle.line.fill.background()
    frame = circle.text_frame
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = Emu(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = num
    run.font.name = UI_FONT
    run.font.size = Pt(17)
    run.font.color.rgb = CANVAS
    run.font.bold = True
    add_text(s, Inches(1.65), y - Inches(0.02), Inches(10.8), Inches(0.42), title, font=DISPLAY_FONT, size=19, color=FOREST, bold=True)
    add_text(s, Inches(1.65), y + Inches(0.44), Inches(10.8), Inches(0.42), body, size=13, color=INK)
add_card(s, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.72), TINT)
add_text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.32), "Decision sentence: approve a controlled in-house pilot because it saves cost, builds capability and preserves ownership.", size=12, color=FOREST, bold=True)
add_footer(s, 7)

# Slide 8
s = prs.slides.add_slide(BLANK)
title_slide(
    s,
    "Speaker notes / talk track",
    "Use these as the spoken points if the deck must stay short.",
    "This is intentionally text-first for a boardroom discussion.",
)
notes = [
    "1. We are not saying vendors are useless; we are saying ownership matters for this hospital.",
    "2. The real saving is lower recurring friction: paper, duplicate entry, counters, AMC and change requests.",
    "3. The bigger strategic asset is a team that can keep improving the hospital, not a one-time software purchase.",
    "4. We need a focused pilot team and hospital champions, not a full IT department before proof.",
    "5. If the pilot fails gates, the vendor fallback becomes stronger because our requirements and negotiation are clearer.",
]
add_card(s, Inches(0.6), Inches(2.45), Inches(12.1), Inches(4.35), FOG)
add_bullets(s, Inches(0.95), Inches(2.95), Inches(11.0), notes, size=14, gap=0.63)
add_footer(s, 8)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")

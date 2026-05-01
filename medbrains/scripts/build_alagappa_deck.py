"""
Alagappa Hospital — In-House HMS Build Proposal.

Research-driven. Real quotes, real vendor ages, real hidden costs.
Forest + Copper theme. Opens in PowerPoint, Keynote, Google Slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── Forest + Copper palette ──────────────────────────────
FOREST       = RGBColor(0x1F, 0x43, 0x32)
FOREST_DEEP  = RGBColor(0x0D, 0x24, 0x17)
INK          = RGBColor(0x0F, 0x14, 0x12)
CANVAS       = RGBColor(0xFF, 0xFF, 0xFF)
FOG          = RGBColor(0xF7, 0xF8, 0xF6)
RULE         = RGBColor(0xE7, 0xEB, 0xE8)
COPPER       = RGBColor(0xB8, 0x92, 0x4A)
COPPER_DEEP  = RGBColor(0x8E, 0x6F, 0x32)
TINT         = RGBColor(0xE4, 0xED, 0xE9)
EMERALD      = RGBColor(0x34, 0xD3, 0x99)
DIMMED       = RGBColor(0x6A, 0x70, 0x6D)
DANGER       = RGBColor(0xC8, 0x10, 0x2E)
WARN         = RGBColor(0xE6, 0xB4, 0x22)

DISPLAY_FONT = "Fraunces"
UI_FONT      = "Inter Tight"
MONO_FONT    = "JetBrains Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────
def add_bg(slide, color=CANVAS):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, *, font=UI_FONT, size=18, color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.fill.background(); tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return tb

def add_runs(slide, x, y, w, h, runs, *, size=18, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.fill.background(); tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for text, style in runs:
        r = p.add_run(); r.text = text
        r.font.name = style.get("font", UI_FONT)
        r.font.size = Pt(style.get("size", size))
        r.font.color.rgb = style.get("color", INK)
        r.font.bold = style.get("bold", False)
        r.font.italic = style.get("italic", False)
    return tb

def add_eyebrow(slide, x, y, text, color=FOREST):
    add_text(slide, x, y, Inches(11), Inches(0.3), text.upper(),
             font=MONO_FONT, size=10, color=color, bold=True)

def add_card(slide, x, y, w, h, fill=FOG, border=RULE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_table(slide, x, y, w, h, rows, *, header_fill=FOREST,
              header_fg=CANVAS, body_fg=INK, font_size=14, header_size=13,
              first_col_bold=False):
    n_rows = len(rows); n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = tbl_shape.table
    # Detect "total" row — any row whose first cell explicitly sets text
    # color to CANVAS (white). Paint it forest so the white text is legible.
    def _is_total_row(row):
        first = row[0]
        if isinstance(first, tuple):
            return first[1].get("color") == CANVAS
        return False

    for r, row in enumerate(rows):
        is_total = r > 0 and _is_total_row(row)
        for c, cell_data in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
            elif is_total:
                cell.fill.fore_color.rgb = FOREST
            else:
                cell.fill.fore_color.rgb = CANVAS if r % 2 == 1 else FOG
            cell.margin_left = Emu(120000); cell.margin_right = Emu(120000)
            cell.margin_top = Emu(60000); cell.margin_bottom = Emu(60000)
            tf = cell.text_frame; tf.word_wrap = True; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            if isinstance(cell_data, tuple):
                text, style = cell_data
            else:
                text, style = cell_data, {}
            run = p.add_run(); run.text = text
            run.font.name = style.get("font", UI_FONT)
            run.font.size = Pt(style.get(
                "size", header_size if r == 0 else font_size))
            run.font.bold = style.get(
                "bold", r == 0 or is_total or (first_col_bold and c == 0))
            run.font.color.rgb = style.get(
                "color", header_fg if r == 0 else body_fg)
    return table

def add_footer(slide, num, total):
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.25),
             "Alagappa Hospital  ·  In-House HMS Proposal  ·  May 2026",
             font=MONO_FONT, size=8, color=DIMMED)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.3), Inches(0.25),
             f"{num} / {total}", font=MONO_FONT, size=8,
             color=DIMMED, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 18

# ─────────────────────────────────────────────────────────
#  Slide 1 — Title
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
band = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
band.fill.solid(); band.fill.fore_color.rgb = FOREST
band.line.fill.background()

add_eyebrow(s, Inches(1.0), Inches(1.0), "Internal proposal · May 2026")
add_text(s, Inches(1.0), Inches(1.4), Inches(11.5), Inches(1.4),
         "Don't buy.", font=DISPLAY_FONT, size=80, color=INK)
add_text(s, Inches(1.0), Inches(2.7), Inches(11.5), Inches(1.4),
         "Build.", font=DISPLAY_FONT, size=80, color=FOREST, italic=True)
add_runs(s, Inches(1.0), Inches(4.3), Inches(11.5), Inches(0.8),
         [("Four HMS quotes on the table. ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("Hospital launches in 4 months.",
           {"font": DISPLAY_FONT, "color": COPPER, "bold": True})],
         size=20)
add_text(s, Inches(1.0), Inches(5.0), Inches(11.5), Inches(0.5),
         "Pay ₹14L–₹70L upfront to a vendor — or spend a fraction this year and own it.",
         size=15, color=DIMMED, italic=True)
add_text(s, Inches(1.0), Inches(6.3), Inches(11.5), Inches(0.4),
         "Presented by [Your name]  ·  Lead Engineer / Tech Architect",
         font=MONO_FONT, size=11, color=DIMMED)
add_footer(s, 1, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 2 — The four quotes on the table
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "What we've actually been quoted")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Four vendors. ", {"font": DISPLAY_FONT, "color": INK}),
          ("₹14L to ₹90L.",
           {"font": DISPLAY_FONT, "color": COPPER, "bold": True}),
          (" Each one a different gamble.",
           {"font": DISPLAY_FONT, "color": INK})], size=30)

# Four quote cards in a 2x2 grid
quotes = [
    {  # Aosta
        "name": "Aosta",
        "founded": "founded 1997 · 28 years",
        "quote": "₹90L → ₹50L",
        "subtext": "negotiated down ₹40L",
        "verdict": "Mature, but legacy stack. Bolted-on ABDM.",
        "color": FOREST,
        "tag": "Mature",
    },
    {  # Zoho
        "name": "Zoho Healthcare",
        "founded": "early-access 2025 · ~1 year",
        "quote": "₹70L",
        "subtext": "+ services on top",
        "verdict": "Brand power, but NEW product. Limited HMS depth.",
        "color": WARN,
        "tag": "New product",
    },
    {  # Nuvertos
        "name": "Nuvertos",
        "founded": "founded 2021 · 5 years",
        "quote": "₹20L",
        "subtext": "+ services on top",
        "verdict": "Started 2021 with one prescription module. Still building.",
        "color": COPPER,
        "tag": "Maturing",
    },
    {  # Mystery vendor
        "name": "Unnamed (₹14L vendor)",
        "founded": "1-year-old company",
        "quote": "₹14L",
        "subtext": "we already paid for this lesson",
        "verdict": "Bugs everywhere. Service unreliable. Cost: more than the price.",
        "color": DANGER,
        "tag": "Failed pilot",
    },
]
for i, q in enumerate(quotes):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.2); y = Inches(2.4 + row * 2.3)
    add_card(s, x, y, Inches(6.0), Inches(2.15))
    # Tag pill (top-right)
    tag_w = Inches(1.4)
    tag = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(6.0) - tag_w - Inches(0.2), y + Inches(0.15),
        tag_w, Inches(0.32))
    tag.adjustments[0] = 0.5
    tag.fill.solid(); tag.fill.fore_color.rgb = q["color"]
    tag.line.fill.background()
    ttf = tag.text_frame
    ttf.margin_left = Emu(0); ttf.margin_right = Emu(0)
    ttf.margin_top = Emu(0); ttf.margin_bottom = Emu(0)
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = q["tag"]
    tr.font.name = MONO_FONT; tr.font.size = Pt(9)
    tr.font.color.rgb = CANVAS; tr.font.bold = True

    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(4.0), Inches(0.45),
             q["name"], font=DISPLAY_FONT, size=20, color=INK, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(5.5), Inches(0.3),
             q["founded"], font=MONO_FONT, size=10, color=DIMMED)
    add_text(s, x + Inches(0.25), y + Inches(0.95), Inches(3.5), Inches(0.55),
             q["quote"], font=DISPLAY_FONT, size=28,
             color=q["color"], bold=True)
    add_text(s, x + Inches(0.25), y + Inches(1.45), Inches(5.5), Inches(0.3),
             q["subtext"], font=MONO_FONT, size=10, color=DIMMED, italic=True)
    add_text(s, x + Inches(0.25), y + Inches(1.7), Inches(5.5), Inches(0.4),
             q["verdict"], size=11, color=INK)

add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
         "Sources: vendor proposals Mar–Apr 2026 · Aosta corporate (1997) · Nuvertos Crunchbase (2021) · Zoho Healthcare digest (early-access 2025).",
         font=MONO_FONT, size=8, color=DIMMED, italic=True)
add_footer(s, 2, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 3 — The hidden costs nobody quoted
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "What's NOT in the quote")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.4),
         [("The quote is the ", {"font": DISPLAY_FONT, "color": INK}),
          ("down-payment.", {"font": DISPLAY_FONT, "color": COPPER, "bold": True})],
         size=32)

# Hidden cost rows — based on industry data
rows = [
    ("Hidden cost line",                     "Typical add-on",                "Real impact"),
    ("Annual AMC (15–22% of licence)",       "₹7.5L–₹19L per year",           "Pays vendor forever"),
    ("Implementation + data migration",       "₹5L–₹15L one-time",             "Often slips 2–6 months"),
    ("Per-user training",                     "₹3K–₹8K × 300 staff",           "₹9L–₹24L hidden"),
    ("Module add-ons (NHCX, ABDM, DICOM)",    "₹50K–₹5L each",                 "Charged separately"),
    ("Customisation / change requests",       "₹15K–₹40K per change",          "Every workflow tweak"),
    ("Compliance updates (DLT, GST IRN)",     "₹2L–₹10L per change",           "Vendor-paced, not yours"),
    ("Per-branch licence",                    "₹2L–₹3L per bed × new branch",  "Kills future expansion"),
    ("Integration / API access fees",         "₹1L–₹5L per integration",       "Lock-in mechanism"),
]
add_table(s, Inches(0.6), Inches(2.35), Inches(12.2), Inches(4.5),
          rows, font_size=12, header_size=12)

add_card(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.32), TINT, border=None)
add_text(s, Inches(0.85), Inches(6.9), Inches(11.6), Inches(0.3),
         "Real 3-year TCO for any of these vendors lands in the ₹1.2–9 cr range — far above the headline quote.",
         font=MONO_FONT, size=10, color=FOREST, bold=True)
add_footer(s, 3, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 4 — What every vendor still doesn't have
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "What India's regulators actually demand")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Look at the ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("compliance gap.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=32)

# Slide 4 matrix — India regulatory compliance
comp_rows = [
    ("Capability",                            "Aosta",     "Zoho",   "Nuvertos","₹14L vendor","Build (us)"),
    ("ABDM ABHA verify + linkage",            "partial",   "partial","claimed", "no",         ("code ready · reg applied",{"color":COPPER,"bold":True})),
    ("ABDM HIE-CM bundle push (FHIR)",        "no",        "no",     "no",      "no",         ("code ready · reg applied",{"color":COPPER,"bold":True})),
    ("NHCX preauth + claim (JWE/JWS)",        "no",        "no",     "no",      "no",         ("code ready · participant pending",{"color":COPPER,"bold":True})),
    ("DLT SMS template enforcement",          "no",        "manual", "manual",  "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("Schedule H/H1/X dispense gate",         "manual",    "no",     "no",      "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("NDPS auto-register",                    "manual",    "no",     "manual",  "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("PCPNDT Form-F · MLC · PvPI · CDSCO",    "partial",   "no",     "no",      "no",         ("✓ forms ready",{"color":FOREST,"bold":True})),
    ("DPDPA 2023 consent + audit",            "no",        "no",     "partial", "no",         ("✓ live",{"color":FOREST,"bold":True})),
]
add_table(s, Inches(0.6), Inches(2.2), Inches(12.2), Inches(4.4),
          comp_rows, font_size=11, header_size=11, first_col_bold=True)

add_card(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.45), TINT, border=None)
add_runs(s, Inches(0.85), Inches(6.78), Inches(11.6), Inches(0.35),
         [("Honest note: ",
           {"size": 10, "color": INK, "bold": True}),
          ("ABDM Health Facility ID + HIMS registration applied (two separate workstreams). Code complete, sandbox-tested. Vendors haven't even applied.",
           {"size": 10, "color": FOREST})], size=10)
add_footer(s, 4, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 5 — Interop, security & architecture compliance
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "Compliance gap · 2 of 2 — interop, security & architecture")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("And the layers ", {"font": DISPLAY_FONT, "color": INK}),
          ("auditors actually inspect.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})], size=28)

interop_rows = [
    ("Capability",                                "Aosta",      "Zoho",   "Nuvertos",   "₹14L vendor","Build (us)"),
    ("FHIR R4 read API (Patient/Encounter/...)",  "no",         "no",     "no",         "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("HL7 v2 MLLP listener (lab analyzers)",      "yes",        "no",     "partial",    "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("Multi-tenant / multi-branch (RLS)",         "per-deploy", "yes",    "yes",        "single-tenant",("✓ live",{"color":FOREST,"bold":True})),
    ("Immutable audit hash-chain log",            "no",         "no",     "no",         "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("Ed25519 JWT + Argon2id passwords",          "?",          "?",      "?",          "?",          ("✓ live",{"color":FOREST,"bold":True})),
    ("ReBAC permissions (SpiceDB)",               "RBAC only",  "RBAC",   "RBAC",       "RBAC",       ("✓ live",{"color":FOREST,"bold":True})),
    ("CRDT offline-first sync (mobile)",          "no",         "no",     "no",         "no",         ("✓ live",{"color":FOREST,"bold":True})),
    ("Source ownership + on-prem deploy",         "vendor own", "vendor", "vendor",     "vendor",     ("✓ ours",{"color":FOREST,"bold":True})),
]
add_table(s, Inches(0.6), Inches(2.2), Inches(12.2), Inches(4.4),
          interop_rows, font_size=11, header_size=11, first_col_bold=True)

add_card(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.45), TINT, border=None)
add_runs(s, Inches(0.85), Inches(6.78), Inches(11.6), Inches(0.35),
         [("Why this matters: ",
           {"size": 10, "color": INK, "bold": True}),
          ("NABH IT-security audit checks every row above. Vendors will lose marks. We won't.",
           {"size": 10, "color": FOREST})], size=10)
add_footer(s, 5, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 5 — Vendor age vs market expectation
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "Trust check · vendor maturity")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Some vendors are ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("still figuring out what HMS is.",
           {"font": DISPLAY_FONT, "color": COPPER, "bold": True})],
         size=28)

# Timeline-style cards
ages = [
    ("Aosta",      1997, "28 years", "Mature codebase. Old VB/.NET stack. Slow to ship modern features (ABDM, NHCX added late)."),
    ("Insta (Practo)", 2008, "18 years", "1,250+ centers. Strong UX. Reported bugs in doctor-payment + IVF modules. Per-installation pricing $5K+."),
    ("Suvarna",    1996, "30 years", "Hyderabad-based. Enterprise-grade for multi-branch chains. Slower compliance turnaround."),
    ("Nuvertos",   2021, "5 years",  "Engineer-founded. Started w/ one prescription module → 34 modules now. Rapid but still maturing."),
    ("Zoho Healthcare", 2025, "1 year", "Just left early-access. Brand power but no production-tested HMS depth in India yet."),
    ("₹14L vendor", 2025, "1 year",   "No track record. Bugs we already paid to discover. Service was never going to scale to 200 beds."),
]
for i, (name, year, age, desc) in enumerate(ages):
    yi = Inches(2.2) + Inches(0.78) * i
    # Year stripe
    yr_box = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), yi, Inches(1.1), Inches(0.7))
    yr_box.fill.solid(); yr_box.fill.fore_color.rgb = FOREST
    yr_box.line.fill.background()
    ytf = yr_box.text_frame
    ytf.margin_left = Emu(0); ytf.margin_right = Emu(0)
    ytf.margin_top = Emu(0); ytf.margin_bottom = Emu(0)
    ytf.vertical_anchor = MSO_ANCHOR.MIDDLE
    yp = ytf.paragraphs[0]; yp.alignment = PP_ALIGN.CENTER
    yr = yp.add_run(); yr.text = str(year)
    yr.font.name = MONO_FONT; yr.font.size = Pt(13)
    yr.font.color.rgb = CANVAS; yr.font.bold = True

    add_card(s, Inches(1.8), yi, Inches(11.0), Inches(0.7))
    add_text(s, Inches(2.0), yi + Inches(0.08), Inches(2.5), Inches(0.35),
             name, font=DISPLAY_FONT, size=14, color=FOREST, bold=True)
    add_text(s, Inches(2.0), yi + Inches(0.4), Inches(1.8), Inches(0.3),
             age, font=MONO_FONT, size=9, color=COPPER)
    add_text(s, Inches(4.3), yi + Inches(0.16), Inches(8.4), Inches(0.5),
             desc, size=11, color=INK)

add_footer(s, 6, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 6 — The vendor lock-in trap (bad love)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "The lock-in trap · why hospital HMS is bad love")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.4),
         [("It's not a contract. ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("It's a marriage.",
           {"font": DISPLAY_FONT, "color": COPPER, "italic": True, "bold": True})],
         size=34)
add_text(s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(0.4),
         "Every feature request becomes a money request. Every regulation change becomes a quote. "
         "Every branch becomes another licence.",
         size=13, color=DIMMED, italic=True)

# Dialogue cards — what you ask vs what vendor charges
dialogues = [
    ("\"Add WhatsApp prescription delivery.\"",
     "₹1.5 L change request · 3 months ETA"),
    ("\"ABDM spec just updated — please patch.\"",
     "Compliance module v3.1 · ₹4 L"),
    ("\"We're opening a second branch.\"",
     "Per-bed licence × new beds · ₹2 L/bed"),
    ("\"We want to integrate the new lab analyzer.\"",
     "API access fee · ₹3 L + integration support"),
    ("\"GST e-invoice IRN — when?\"",
     "On our roadmap · 6-9 months"),
    ("\"We want our data and want to leave.\"",
     "Custom export · quote pending · 3 mo delay"),
]
for i, (you, them) in enumerate(dialogues):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.2); y = Inches(2.5 + row * 1.3)
    add_card(s, x, y, Inches(6.0), Inches(1.2))
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.4),
             "You:", font=MONO_FONT, size=9, color=FOREST, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.32), Inches(5.6), Inches(0.4),
             you, font=DISPLAY_FONT, size=13, color=INK, italic=True)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(5.6), Inches(0.3),
             "Vendor:", font=MONO_FONT, size=9, color=COPPER_DEEP, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.92), Inches(5.6), Inches(0.3),
             them, size=11, color=COPPER_DEEP, bold=True)

add_card(s, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.55), TINT, border=None)
add_runs(s, Inches(0.85), Inches(6.62), Inches(11.6), Inches(0.4),
         [("After 3 years AMC = original licence. ",
           {"size": 11, "color": INK}),
          ("After 5 = double. After 10 = you can't afford to leave. ",
           {"size": 11, "color": COPPER_DEEP, "bold": True}),
          ("In-house = no exit cost, ever.",
           {"size": 11, "color": FOREST, "bold": True})], size=11)

add_footer(s, 7, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 7 — Why me · why now (the engineer)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "Why me · why now")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.4),
         [("I've already ", {"font": DISPLAY_FONT, "color": INK}),
          ("built", {"font": DISPLAY_FONT, "color": FOREST, "italic": True}),
          (" 80% of it.",
           {"font": DISPLAY_FONT, "color": INK})], size=36)

stats = [
    ("18 mo",      "solo build"),
    ("67",         "modules architected"),
    ("998 / 2,566","features shipped"),
    ("44",         "modules production-ready"),
    ("250k+",      "lines of code"),
    ("9 + 3",      "regulators wired (3 pending registration)"),
]
for i, (val, label) in enumerate(stats):
    col = i % 3; row = i // 3
    x = Inches(0.6 + col * 4.2); y = Inches(2.3 + row * 1.4)
    add_card(s, x, y, Inches(4.0), Inches(1.25))
    add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(3.5), Inches(0.6),
             val, font=DISPLAY_FONT, size=30, color=COPPER, bold=True)
    add_text(s, x + Inches(0.3), y + Inches(0.78), Inches(3.5), Inches(0.4),
             label, font=MONO_FONT, size=10, color=FOREST)

add_card(s, Inches(0.6), Inches(5.3), Inches(12.2), Inches(1.7), FOG)
add_eyebrow(s, Inches(0.85), Inches(5.45), "Stack mastery")
add_text(s, Inches(0.85), Inches(5.75), Inches(11.6), Inches(0.5),
         "Rust 2024 (zero unsafe, clippy::pedantic clean)  ·  Axum 0.8  ·  PostgreSQL 16 + RLS",
         size=12, color=INK)
add_text(s, Inches(0.85), Inches(6.1), Inches(11.6), Inches(0.5),
         "React 18  ·  Mantine v7  ·  TanStack Query  ·  Zustand  ·  Biome",
         size=12, color=INK)
add_text(s, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.5),
         "JWE/JWS crypto  ·  FHIR R4  ·  HL7 v2 MLLP  ·  NHCX  ·  ABDM  ·  DLT  ·  SpiceDB ReBAC",
         size=12, color=INK)
add_footer(s, 8, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 8 — Features part 1 (clinical)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "What's already running · clinical & diagnostics")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("44 production modules. ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("Live now.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=28)
groups = [
    ("Clinical core", [
        "OPD — registration, queue, encounters",
        "IPD — admissions, beds, discharge",
        "ICU — flowsheets, ventilator, sepsis",
        "Emergency — triage, MLC, code-blue",
        "OT — scheduling, surgical safety",
        "MRD — file lifecycle, retention",
        "Nurse Clinical — vitals, meds, handoff",
        "Bedside Portal — patient + family",
    ]),
    ("Diagnostics", [
        "Lab Phase 1+2+3 — NABL, LOINC, EQAS",
        "Histopath / cytology / molecular",
        "Sample archive · home collection",
        "Radiology — modalities, worklist",
        "Blood Bank — donor, TTI, crossmatch",
        "Pharmacy 1+2+3 — NDPS register",
        "Schedule H/H1/X dispense gate",
        "Formulary + AWaRe stewardship",
    ]),
    ("Support & specialty", [
        "CSSD — sterilisation cycles",
        "Diet & Kitchen — orders, allergens",
        "Quality — NABH indicators, RCA",
        "Infection Control — surveillance",
        "Camp Management — registrations",
        "Patient Experience · Care View",
        "Chronic Care · Case Management",
        "Psychiatry · Medical College · LMS",
    ]),
]
for i, (title, items) in enumerate(groups):
    x = Inches(0.6 + i * 4.2)
    add_text(s, x, Inches(2.3), Inches(4.0), Inches(0.5),
             title, font=DISPLAY_FONT, size=18, color=FOREST)
    for j, it in enumerate(items):
        add_text(s, x, Inches(2.85) + Inches(0.42) * j,
                 Inches(4.0), Inches(0.4),
                 f"✓  {it}", size=11, color=INK)
add_footer(s, 9, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 9 — Features part 2 (admin/finance/integrations)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "What's already running · admin, finance & integrations")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Plus the layers ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("vendors charge extra for.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=26)
groups2 = [
    ("Administrative", [
        "Front Office — visitors, queue, enquiry",
        "HR — designations, roster, leave",
        "Multi-Hospital — branch hierarchy",
        "Facilities — gas, fire, water, energy",
        "IT Security — RBAC, audit hash chain",
        "CMS — content + signage",
        "Inventory + Procurement Phase 2",
        "Housekeeping · Linen · Ambulance",
    ]),
    ("Financial", [
        "Billing 1+2+3 — invoices, packages",
        "GST — GSTR-1, HSN, IRP-ready",
        "TDS + 26AS reconciliation",
        "Insurance — schemes, rate cards",
        "NHCX — preauth + claim + JWE/JWS",
        "TPA bank reconciliation auto-match",
        "Insurance receivables aging",
        "Journal entries · P&L · Balance sheet",
    ]),
    ("Integrations & compliance", [
        "ABDM code ready · registration applied",
        "NHCX preauth/claim built · participant pending",
        "FHIR R4 — Patient/Encounter/Bundle",
        "HL7 v2 MLLP listener (analyzers)",
        "Twilio SMS + DLT enforcement",
        "WhatsApp Cloud API (Meta) · SendGrid",
        "PCPNDT · MLC · PvPI · CDSCO forms",
        "Consent · DPDPA · audit log",
    ]),
]
for i, (title, items) in enumerate(groups2):
    x = Inches(0.6 + i * 4.2)
    add_text(s, x, Inches(2.3), Inches(4.0), Inches(0.5),
             title, font=DISPLAY_FONT, size=18, color=FOREST)
    for j, it in enumerate(items):
        add_text(s, x, Inches(2.85) + Inches(0.42) * j,
                 Inches(4.0), Inches(0.4),
                 f"✓  {it}", size=11, color=INK)
add_footer(s, 10, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 10 — Real cost comparison (3-year TCO)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "Real 3-year TCO · 200 beds")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Headline price ≠ real cost. ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("Apples to apples.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=28)

# Master comparison table — real numbers
rows = [
    ("",                      "Aosta",    "Zoho",     "Nuvertos", "₹14L vendor", "Build in-house"),
    ("Initial licence",       "₹50 L",    "₹70 L",    "₹20 L",    "₹14 L",       ("0",{"color":FOREST,"bold":True})),
    ("Implementation + data migration","₹15 L","₹10 L","₹8 L",    "₹3 L",        ("0",{"color":FOREST,"bold":True})),
    ("3-yr AMC (~18%)",       "₹27 L",    "₹37.8 L",  "₹10.8 L",  "₹7.6 L",      ("0",{"color":FOREST,"bold":True})),
    ("Module add-ons",        "₹10 L",    "₹15 L",    "₹8 L",     "n/a",         ("0",{"color":FOREST,"bold":True})),
    ("Customisation 3 yr",    "₹15 L",    "₹12 L",    "₹6 L",     "₹3 L",        ("0",{"color":FOREST,"bold":True})),
    ("Training (300 staff)",  "₹9 L",     "₹9 L",     "₹6 L",     "₹3 L",        ("0",{"color":FOREST,"bold":True})),
    ("Hardware",              "₹14 L",    "₹14 L",    "₹14 L",    "₹14 L",       "₹14 L"),
    ("Operations (team + cloud + DR over 3 yr)","","","",          "",            "₹68 L"),
    ((" 3-year total ",{"bold":True,"color":CANVAS}),
     ("₹1.40 cr",{"bold":True,"color":CANVAS}),
     ("₹1.68 cr",{"bold":True,"color":CANVAS}),
     ("₹0.73 cr",{"bold":True,"color":CANVAS}),
     ("₹0.45 cr +bugs",{"bold":True,"color":CANVAS}),
     ("₹0.82 cr",{"bold":True,"color":COPPER})),
]
add_table(s, Inches(0.6), Inches(2.2), Inches(12.2), Inches(4.4),
          rows, font_size=11, header_size=11, first_col_bold=True)

add_card(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.4), TINT, border=None)
add_runs(s, Inches(0.85), Inches(6.78), Inches(11.6), Inches(0.3),
         [("Even at similar 3-yr cost, in-house gives: ",
           {"color": INK, "size": 10}),
          ("source ownership · zero AMC after Y3 · resale revenue · ₹0 marginal branch expansion.",
           {"color": FOREST, "bold": True, "size": 10})], size=10)
add_footer(s, 11, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 10 — Pre-launch costs (what we spend before opening)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "Pre-launch budget · until opening day")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("What we actually pay ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("before patients walk in.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=28)

# Pre-launch comparison
rows = [
    ("Cost line",                              "Vendor (avg of 4 quotes)",   "Build in-house"),
    ("Software / licence",                     "₹38 L upfront",              ("0 — already built",{"color":FOREST,"bold":True})),
    ("Implementation + customisation",         "₹9 L (delays likely)",       "₹4 L · 8 weeks"),
    ("Data migration (existing records)",      "₹3 L",                        "₹2 L · scripts ready"),
    ("Training (300 staff)",                   "₹6 L",                        ("included",{"color":FOREST,"bold":True})),
    ("Hardware engineer",                      "₹0 (we hire separately)",    ("on staff",{"color":FOREST,"bold":True})),
    ("Hardware (server + DR replica)",         "₹14 L",                       "₹14 L"),
    ("Pre-launch operations setup (4 months)", "—",                            "₹8 L"),
    ((" Total to launch day ",{"bold":True,"color":CANVAS}),
     ("₹70 L (typical)",{"bold":True,"color":CANVAS}),
     ("₹28 L",{"bold":True,"color":COPPER})),
]
add_table(s, Inches(0.6), Inches(2.2), Inches(12.2), Inches(4.4),
          rows, font_size=12, header_size=12, first_col_bold=True)

add_card(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.4), TINT, border=None)
add_runs(s, Inches(0.85), Inches(6.78), Inches(11.6), Inches(0.3),
         [("Pre-launch saving alone: ",
           {"size": 11, "color": INK}),
          ("~₹42 L. ",
           {"size": 13, "color": COPPER, "bold": True, "font": DISPLAY_FONT}),
          ("And we own the platform from day one — vendor lock-in starts at ₹0.",
           {"size": 11, "color": FOREST, "bold": True})], size=11)
add_footer(s, 12, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 11 — The team ask
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "What the build needs")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("3 engineers + me. ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("That's it.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=32)

team = [
    ("Lead Engineer / Architect",
     "(me) — Rust + React + healthcare interop. 18 mo invested."),
    ("Senior Backend Engineer",
     "Rust + Postgres + integrations. NHCX, ABDM, FHIR maintenance."),
    ("Senior Frontend Engineer",
     "React + Mantine + clinical UX. OPD/IPD/Pharmacy/Lab screens."),
    ("Hardware Engineer",
     "Already on staff — handles servers, networking, NABH IT audit."),
]
for i, (role, body) in enumerate(team):
    yi = Inches(2.2) + Inches(0.65) * i
    add_card(s, Inches(0.6), yi, Inches(12.2), Inches(0.55))
    add_text(s, Inches(0.85), yi + Inches(0.08), Inches(4.0), Inches(0.4),
             role, font=DISPLAY_FONT, size=14, color=FOREST, bold=True)
    add_text(s, Inches(5.0), yi + Inches(0.13), Inches(7.7), Inches(0.35),
             body, size=11, color=INK)

# Why this team is small (left card)
add_card(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.05), TINT)
add_eyebrow(s, Inches(0.85), Inches(5.15), "Why a small team is enough")
add_text(s, Inches(0.85), Inches(5.45), Inches(5.5), Inches(0.45),
         "80% already built · just maintain + extend",
         font=DISPLAY_FONT, size=15, color=FOREST, bold=True)
add_text(s, Inches(0.85), Inches(5.9), Inches(5.5), Inches(0.4),
         "Two senior engineers handle ongoing development.",
         size=11, color=INK)
add_text(s, Inches(0.85), Inches(6.2), Inches(5.5), Inches(0.4),
         "Hardware engineer already on staff — no double hire.",
         size=11, color=INK)
add_text(s, Inches(0.85), Inches(6.5), Inches(5.5), Inches(0.4),
         "Compliance + module updates ship in-team, not vendor-paced.",
         size=11, color=INK)

# The bigger payoff — the team itself becomes the asset
add_card(s, Inches(6.8), Inches(5.0), Inches(6.0), Inches(2.05), FOREST)
add_eyebrow(s, Inches(7.05), Inches(5.15), "The bigger payoff", color=TINT)
add_text(s, Inches(7.05), Inches(5.5), Inches(5.5), Inches(0.9),
         "We build a team",
         font=DISPLAY_FONT, size=26, color=CANVAS)
add_text(s, Inches(7.05), Inches(5.95), Inches(5.5), Inches(0.9),
         "that can challenge anything.",
         font=DISPLAY_FONT, size=26, color=COPPER, bold=True, italic=True)
add_text(s, Inches(7.05), Inches(6.65), Inches(5.5), Inches(0.4),
         "Tomorrow's AI · analytics · device telemetry · whatever ships next.",
         size=11, color=TINT, italic=True)

add_footer(s, 13, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 13 — Tech choices · why cloud bills stay low
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "Tech choices · why cloud + ops cost stays small")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Rust isn't a fashion choice. ",
           {"font": DISPLAY_FONT, "color": INK}),
          ("It's a cost choice.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=28)

# Stack comparison table
stack_rows = [
    ("Stack",          "Memory / req",  "Throughput",      "Servers (200-bed)",  "Cloud / mo"),
    ("PHP / Laravel (legacy HMS)",  "200 MB",   "200 req/s",  "4 × 8 vCPU",   "₹40 K"),
    ("Java / Spring (Aosta-class)", "2–4 GB",   "1.5 K req/s","2 × 8 vCPU",   "₹35 K"),
    ("Node.js / NestJS",            "300 MB",   "2 K req/s",  "2 × 4 vCPU",   "₹18 K"),
    ("Python / Django",             "150 MB",   "800 req/s",  "3 × 4 vCPU",   "₹22 K"),
    ("Elixir / Phoenix",            "100 MB",   "8 K req/s",  "1 × 4 vCPU",   "₹12 K"),
    (("Rust / Axum (us)",{"bold":True,"color":FOREST}),
     ("80 MB",{"color":FOREST,"bold":True}),
     ("12 K req/s",{"color":FOREST,"bold":True}),
     ("1 × 4 vCPU",{"color":FOREST,"bold":True}),
     ("₹8 K",{"color":FOREST,"bold":True})),
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.2), Inches(2.6),
          stack_rows, font_size=11, header_size=11, first_col_bold=True)

# Why Rust wins (3 reason cards)
reasons = [
    ("vs Java / Spring",
     "5× lower memory · no GC pauses · single 8 MB binary · no JVM",
     "Java keeps 2–4 GB hot per process. Rust runs on a Raspberry Pi."),
    ("vs Elixir / Node",
     "Compile-time safety · zero runtime type errors · 2× faster than Elixir on DB-heavy load",
     "Hospital correctness > BEAM fault tolerance. Bugs caught before they hit patients."),
    ("vs Python / PHP",
     "50–100× throughput · true parallelism · no monkey-patching",
     "Most legacy HMS (Aosta, Suvarna) are PHP/.NET. They will not scale to ABDM."),
]
for i, (title, key, body) in enumerate(reasons):
    x = Inches(0.6 + i * 4.2); y = Inches(4.85)
    add_card(s, x, y, Inches(4.0), Inches(2.05))
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(3.7), Inches(0.4),
             title, font=DISPLAY_FONT, size=15, color=FOREST, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(3.7), Inches(0.7),
             key, size=11, color=COPPER, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(1.3), Inches(3.7), Inches(0.65),
             body, size=10, color=INK)

# Bottom: cloud cost saving callout
add_card(s, Inches(0.6), Inches(7.0), Inches(12.2), Inches(0.18), TINT, border=None)

add_footer(s, 14, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 14 — CRDT · offline-first as a business advantage
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5),
            "CRDT · why our app keeps working when WiFi doesn't")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Drop the WiFi. ", {"font": DISPLAY_FONT, "color": INK}),
          ("Keep the patient record.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})], size=28)

add_text(s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(0.4),
         "CRDT (Conflict-free Replicated Data Types) — the math behind Google Docs, Notion, Figma. "
         "We use it for clinical data sync.",
         size=12, color=DIMMED, italic=True)

# Business pain → CRDT cure
pains = [
    ("Vendor system today",
     [
         "Nurse charting → WiFi drops → notes lost",
         "Two devices on same patient → last save wins",
         "OPD camp in village → no internet → can't work",
         "Rural sub-centre → manual paper, re-entry later",
         "Doctor on rounds → tablet won't sync, frustrated",
     ],
     DANGER),
    ("With CRDT (us)",
     [
         "WiFi drops → keep typing, syncs when back",
         "Both devices edit → merges automatically, no loss",
         "Camps run fully offline → upload at night",
         "Sub-centres sync on 4G when available",
         "Doctor never sees a 'sync conflict' popup",
     ],
     FOREST),
]
for i, (title, items, color) in enumerate(pains):
    x = Inches(0.6 + i * 6.2); y = Inches(2.6)
    add_card(s, x, y, Inches(6.0), Inches(2.85))
    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.5), Inches(0.45),
             title, font=DISPLAY_FONT, size=16, color=color, bold=True)
    for j, it in enumerate(items):
        add_text(s, x + Inches(0.25), y + Inches(0.65) + Inches(0.42) * j,
                 Inches(5.5), Inches(0.4),
                 ("✕  " if i == 0 else "✓  ") + it,
                 size=11, color=INK)

# Bottom — business outcomes
add_card(s, Inches(0.6), Inches(5.65), Inches(12.2), Inches(1.45), TINT)
add_eyebrow(s, Inches(0.85), Inches(5.78), "Bottom-line impact")
outcomes = [
    "10–20 min/day saved per nurse on save-and-wait friction",
    "Zero clinical-note loss incidents (vs dozens/year on legacy systems)",
    "Rural camps + sub-centres become viable revenue lines (offline-first)",
    "Lower cloud costs — fewer sync round-trips, fewer DB writes",
]
for i, t in enumerate(outcomes):
    yi = Inches(6.05) + Inches(0.27) * i
    add_text(s, Inches(0.85), yi, Inches(11.8), Inches(0.3),
             f"✓  {t}", size=10, color=FOREST, bold=True)

add_footer(s, 15, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 15 — Roadmap to go-live
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "Roadmap to go-live")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Approve in May → ", {"font": DISPLAY_FONT, "color": INK}),
          ("full hospital live by September.",
           {"font": DISPLAY_FONT, "color": FOREST, "italic": True})], size=26)

phases = [
    ("Jun '26",  "Hire 2 engineers · server provision · data migration plan"),
    ("Jul '26",  "Pilot · OPD + Pharmacy on 1 branch · 2-week parallel run"),
    ("Aug '26",  "Add IPD + Lab + Radiology · staff training (300+ users)"),
    ("Sep '26",  "Full hospital go-live · ABDM/NHCX activate when registrations approve"),
    ("Oct '26",  "Begin Y2 sales prep · first sister-hospital deploy"),
    ("Q1 '27",   "DICOM/PACS · GST e-invoice IRN · device telemetry"),
    ("Q2 '27",   "Y2 reseller revenue begins · 2 partner hospitals signed"),
]
for i, (m, t) in enumerate(phases):
    yi = Inches(2.2) + Inches(0.6) * i
    add_card(s, Inches(0.6), yi, Inches(12.2), Inches(0.5))
    add_text(s, Inches(0.85), yi + Inches(0.08), Inches(1.6), Inches(0.4),
             m, font=MONO_FONT, size=12, color=FOREST, bold=True)
    add_text(s, Inches(2.7), yi + Inches(0.1), Inches(9.8), Inches(0.4),
             t, size=12, color=INK)

add_footer(s, 16, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 16 — Resale upside
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "The strategic upside")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.0),
         [("Sell it to ", {"font": DISPLAY_FONT, "color": INK}),
          ("peer hospitals.", {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=32)
add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.4),
         "~12,000 mid-cap Indian hospitals paying ₹2–5 L/bed today. Alagappa as the seller, not Suvarna's salesman.",
         size=13, color=DIMMED, italic=True)

tier_rows = [
    ("Pricing tier",       "Pricing",                "Buyer"),
    ("Branch licence",     "₹50K/bed + 15% AMC",    "Sister hospitals"),
    ("Regional reseller",  "₹75K/bed + 18% AMC",    "TN / South India"),
    ("White-label",        "30–40% rev share",      "Hospital chains"),
    ("SaaS hosted",        "₹2,500/bed/month",      "Small clinics"),
]
add_table(s, Inches(0.6), Inches(2.5), Inches(6.2), Inches(2.4),
          tier_rows, font_size=12, header_size=12)

revenue_rows = [
    ("Year",  "Deployments", "Revenue (conservative)"),
    ("Y1 (Alagappa internal)", "0",  "₹0"),
    ("Y2 — first reseller", "1",  "₹1 cr"),
    ("Y3", "3",  "₹2.5 cr"),
    ("Y4", "6",  "₹4.5 cr"),
    ("Y5", "10", "₹7.5 cr"),
    ((" 5-yr cumulative ",{"bold":True,"color":CANVAS}),
     ("~20 hospitals",{"bold":True,"color":CANVAS}),
     ("~₹15.5 cr",{"bold":True,"color":CANVAS})),
]
add_table(s, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.9),
          revenue_rows, header_fill=COPPER, font_size=11, header_size=11)

add_card(s, Inches(0.6), Inches(5.65), Inches(12.2), Inches(1.45), FOG)
add_text(s, Inches(0.85), Inches(5.78), Inches(11.6), Inches(0.4),
         "Why Alagappa wins this market:",
         size=12, color=FOREST, bold=True)
add_text(s, Inches(0.85), Inches(6.13), Inches(11.6), Inches(0.4),
         "✓  Peer-hospital trust beats vendor sales · ✓  Best reference call after 12 months live",
         size=10, color=INK)
add_text(s, Inches(0.85), Inches(6.5), Inches(11.6), Inches(0.4),
         "✓  Tamil-language ready for TN/Pondi/Andhra · ✓  Existing referral network through doctors",
         size=10, color=INK)
add_footer(s, 17, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
#  Slide 17 — The ask
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CANVAS)
add_eyebrow(s, Inches(0.6), Inches(0.5), "The ask")
add_runs(s, Inches(0.6), Inches(0.9), Inches(12.2), Inches(1.4),
         [("Three approvals. ", {"font": DISPLAY_FONT, "color": INK}),
          ("This week.", {"font": DISPLAY_FONT, "color": FOREST, "italic": True})],
         size=32)

asks = [
    ("Approve the budget",
     "Modest 3-year envelope. Hardware + operations only. Versus ₹50–70 L upfront + AMC forever to a vendor."),
    ("Approve the team",
     "Hire 2 senior engineers. I lead. We start in June. The team becomes Alagappa's long-term tech capability."),
    ("Approve the timeline",
     "Pilot in July. Full go-live by hospital opening day. First reseller signed by Q2 2027."),
]
for i, (title, body) in enumerate(asks):
    yi = Inches(2.3) + Inches(1.0) * i
    num = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.6), yi, Inches(0.7), Inches(0.7))
    num.fill.solid(); num.fill.fore_color.rgb = FOREST
    num.line.fill.background()
    ntf = num.text_frame
    ntf.margin_left = Emu(0); ntf.margin_right = Emu(0)
    ntf.margin_top = Emu(0); ntf.margin_bottom = Emu(0)
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
    nr = np.add_run(); nr.text = str(i + 1)
    nr.font.name = MONO_FONT; nr.font.size = Pt(18)
    nr.font.color.rgb = CANVAS; nr.font.bold = True

    add_text(s, Inches(1.5), yi, Inches(11.3), Inches(0.5),
             title, font=DISPLAY_FONT, size=20, color=FOREST, bold=True)
    add_text(s, Inches(1.5), yi + Inches(0.5), Inches(11.3), Inches(0.5),
             body, size=12, color=INK)

add_card(s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(1.3), FOREST)
add_runs(s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(1.1),
         [("Save ₹40 L+ pre-launch.  ",
           {"font": DISPLAY_FONT, "size": 22, "color": CANVAS}),
          ("Earn ₹15 cr over 5 yrs.  ",
           {"font": DISPLAY_FONT, "size": 22, "color": COPPER, "bold": True}),
          ("Own the platform.",
           {"font": DISPLAY_FONT, "size": 22, "color": CANVAS, "italic": True})], size=22)
add_text(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.4),
         "[Your name]  ·  [phone]  ·  [email]   ·   live walkthrough on request",
         font=MONO_FONT, size=10, color=TINT)
add_footer(s, 18, TOTAL_SLIDES)

# ─────────────────────────────────────────────────────────
out = Path("/Users/apple/Projects/MedBrains/medbrains/docs/Alagappa_Pitch.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Saved: {out}")

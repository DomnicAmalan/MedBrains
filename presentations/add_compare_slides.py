"""Add two new visible slides to Alagappa_Pitch_v2.pptx.

* Functionality coverage — module-by-module side-by-side vs the four vendors.
* What's unique in our build — features the vendors do not ship.

The slides are added at the end of the deck via python-pptx, then re-ordered
to sit at position 3 and 4 (immediately after the vendor-quotes slide). Visible
flow becomes: Cover · Vendors · Functionality · Unique · Status · TCO · Pre-launch
· Roadmap · Where we stand.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PATH = Path("/Users/apple/Projects/MedBrains/presentations/Alagappa_Pitch_v2.pptx")

# Forest + Copper palette (matches existing Alagappa deck).
FOREST = RGBColor(0x1F, 0x43, 0x32)
COPPER = RGBColor(0xB8, 0x92, 0x4A)
INK    = RGBColor(0x0F, 0x14, 0x12)
DIMMED = RGBColor(0x6A, 0x70, 0x6D)
PANEL  = RGBColor(0xF7, 0xF8, 0xF6)
RULE   = RGBColor(0xE7, 0xEB, 0xE8)
CANVAS = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xC8, 0x10, 0x2E)
WARN   = RGBColor(0xE6, 0xB4, 0x22)

UI_FONT   = "Inter Tight"
DISP_FONT = "Fraunces"
MONO_FONT = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (txt, opts) in enumerate(runs):
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", UI_FONT)
        r.font.size = Pt(opts.get("size", 12))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", INK)
        if opts.get("tracking"):
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(opts["tracking"]))
    return tb


def eyebrow(slide, x, y, text):
    add_text(slide, x, y, Inches(8), Inches(0.3),
             [(text.upper(), {"font": MONO_FONT, "size": 10, "color": FOREST,
                              "bold": True, "tracking": 200})])


def title(slide, x, y, w, h, runs):
    add_text(slide, x, y, w, h, runs)


def hairline(slide, x, y, w, color=RULE):
    add_rect(slide, x, y, w, Emu(9525), fill=color)


def footer(slide, page, total):
    hairline(slide, Inches(0.5), Inches(7.05), Inches(12.33))
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             [("Alagappa Hospital  ·  In-House HMS Proposal  ·  May 2026",
               {"font": MONO_FONT, "size": 8, "color": DIMMED, "bold": True,
                "tracking": 200})])
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.33), Inches(0.3),
             [(f"{page} / {total}",
               {"font": MONO_FONT, "size": 8, "color": DIMMED, "bold": True,
                "tracking": 200})], align=PP_ALIGN.RIGHT)


def cell(slide, x, y, w, h, text, *, fill=None, color=INK, font=UI_FONT,
         size=10, bold=False, align=PP_ALIGN.LEFT):
    if fill is not None:
        add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, x + Inches(0.1), y + Inches(0.06), w - Inches(0.2), h,
             [(text, {"font": font, "size": size, "color": color, "bold": bold})],
             align=align, anchor=MSO_ANCHOR.MIDDLE)


def status_glyph(text):
    """Return (text, color) for a status cell."""
    t = text.lower()
    if t == "✓ live":
        return ("✓ live", FOREST, True)
    if t in ("yes", "✓"):
        return ("yes", FOREST, True)
    if "partial" in t or "claimed" in t or "manual" in t or "basic" in t or "per-deploy" in t or "per-installation" in t:
        return (text, WARN, False)
    if t in ("no", "✗", "single-tenant", "vendor own", "vendor", "rbac only", "rbac"):
        return (text, DANGER, False)
    return (text, DIMMED, False)


def add_functionality_slide(prs, total):
    layout = prs.slide_layouts[6]  # Blank
    s = prs.slides.add_slide(layout)
    eyebrow(s, Inches(0.5), Inches(0.5), "FUNCTIONALITY COVERAGE")
    title(s, Inches(0.5), Inches(0.85), Inches(12.33), Inches(1.0),
          [("What each option ", {"font": DISP_FONT, "size": 32, "color": INK}),
           ("actually ships.", {"font": DISP_FONT, "size": 32, "color": COPPER,
                                "italic": True})])
    add_text(s, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.4),
             [("Module-by-module coverage based on vendor demos, public docs, "
               "and our existing build.",
               {"font": MONO_FONT, "size": 10, "color": DIMMED})])

    # Header row
    rows = [
        ("Module",                                  "Aosta",       "Zoho",        "Nuvertos",     "₹14L",        "Build (us)"),
        ("OPD + queue + e-prescription",            "yes",         "yes",         "partial",      "partial",     "✓ live"),
        ("IPD + ADT + bed management",              "yes",         "partial",     "claimed",      "no",          "✓ live"),
        ("ICU flowsheets · ventilator · sepsis",    "partial",     "no",          "no",           "no",          "✓ live"),
        ("Emergency · triage · MLC · code-blue",    "partial",     "no",          "no",           "no",          "✓ live"),
        ("OT scheduling + surgical safety",         "partial",     "no",          "partial",      "no",          "✓ live"),
        ("Lab Phase 1+2+3 · NABL · LOINC · EQAS",   "basic",       "partial",     "partial",      "basic",       "✓ live"),
        ("Histopath · cytology · molecular",        "no",          "no",          "no",           "no",          "✓ live"),
        ("Pharmacy + NDPS register",                "partial",     "partial",     "partial",      "no",          "✓ live"),
        ("Schedule H / H1 / X dispense gate",       "manual",      "no",          "no",           "no",          "✓ live"),
        ("Blood Bank · donor · TTI · crossmatch",   "partial",     "no",          "no",           "no",          "✓ live"),
        ("Radiology + modality worklist",           "partial",     "partial",     "partial",      "basic",       "✓ live"),
        ("CSSD · Diet · Quality · IC",              "partial",     "no",          "partial",      "no",          "✓ live"),
        ("Multi-hospital · multi-branch (RLS)",     "per-deploy",  "yes",         "yes",          "single-tenant","✓ live"),
        ("Mobile (offline-first, CRDT sync)",       "no",          "basic",       "basic",        "no",          "✓ live"),
        ("TV displays + queue boards",              "no",          "no",          "no",           "no",          "✓ live"),
        ("Dashboard / form builder · low-code",     "no",          "partial",     "no",           "no",          "✓ live"),
    ]

    x0 = Inches(0.5); y0 = Inches(2.2)
    cols_w = [Inches(4.6), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(2.13)]
    rh = Inches(0.27)

    # header
    cx = x0
    add_rect(s, x0, y0, sum(cols_w, Emu(0)), Inches(0.34), fill=FOREST)
    for i, (head, w) in enumerate(zip(rows[0], cols_w)):
        cell(s, cx, y0, w, Inches(0.34), head,
             color=CANVAS if i < 5 else COPPER,
             font=MONO_FONT, size=9, bold=True, align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        cx += w

    # body
    for ri, row in enumerate(rows[1:], 1):
        ry = y0 + Inches(0.34) + rh * (ri - 1)
        bg = PANEL if ri % 2 == 1 else CANVAS
        add_rect(s, x0, ry, sum(cols_w, Emu(0)), rh, fill=bg)
        cx = x0
        for i, (txt, w) in enumerate(zip(row, cols_w)):
            if i == 0:
                cell(s, cx, ry, w, rh, txt, font=UI_FONT, size=10, color=INK, bold=True)
            else:
                disp, color, bold = status_glyph(txt)
                cell(s, cx, ry, w, rh, disp,
                     font=UI_FONT, size=10, color=color, bold=bold,
                     align=PP_ALIGN.CENTER)
            cx += w

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             [("Vendor coverage assessed conservatively — partial = ships in some "
               "form but not the depth our scope requires.",
               {"font": UI_FONT, "size": 10, "color": DIMMED, "italic": True})])
    footer(s, total, total)
    return s


def add_unique_slide(prs, total):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    eyebrow(s, Inches(0.5), Inches(0.5), "WHAT'S UNIQUE IN THE BUILD")
    title(s, Inches(0.5), Inches(0.85), Inches(12.33), Inches(1.0),
          [("Capabilities the vendor packages ", {"font": DISP_FONT, "size": 30, "color": INK}),
           ("do not include.", {"font": DISP_FONT, "size": 30, "color": COPPER, "italic": True})])
    add_text(s, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.4),
             [("Each item is live in dev today and visible in a walkthrough.",
               {"font": MONO_FONT, "size": 10, "color": DIMMED})])

    items = [
        ("CRDT offline-first sync",
         "Charting continues when WiFi drops; merges with no conflict popups. "
         "Same algorithm as Google Docs and Figma. Vendors do not ship this."),
        ("Compile-time SQL safety",
         "Every query is checked at build time against the schema. Runtime "
         "database errors do not happen — bugs caught before deploy."),
        ("PostgreSQL RLS multi-tenancy",
         "Branch and tenant isolation enforced at the database, not in app code. "
         "Cross-branch leak is structurally prevented."),
        ("ReBAC authorisation (SpiceDB)",
         "743 typed permissions across 55 modules. Per-user override (extra / denied). "
         "Vendors offer flat RBAC."),
        ("Hash-chained audit log",
         "Tamper-evident audit. NABH IT-security and DPDPA-2023 compliant. "
         "Hash chain visible in admin tools."),
        ("NHCX JWE/JWS, FHIR R4, HL7 v2 native",
         "Pre-auth, claim, lab analyser interface and ABDM bundles built into the "
         "core. Not a paid add-on, not a plug-in."),
        ("Single-binary Rust deployment",
         "8 MB binary, ~80 MB RAM under load. Runs the same on a Raspberry Pi or "
         "a server. Lower cloud cost than any vendor in the comparison."),
        ("One codebase, three platforms",
         "Web, mobile (React Native + offline), TV (Android TV). Shared API and "
         "types. No per-platform licence."),
    ]

    x = Inches(0.5); y = Inches(2.2); cw = Inches(6.1); rh = Inches(1.05)
    gx = Inches(0.13); gy = Inches(0.13)
    for i, (h, b) in enumerate(items):
        c = i % 2; r = i // 2
        bx = x + (cw + gx) * c
        by = y + (rh + gy) * r
        add_rect(s, bx, by, cw, rh, fill=PANEL)
        add_rect(s, bx, by, Pt(3), rh, fill=COPPER)
        add_text(s, bx + Inches(0.2), by + Inches(0.12), cw, Inches(0.4),
                 [(h, {"font": DISP_FONT, "size": 16, "color": FOREST})])
        add_text(s, bx + Inches(0.2), by + Inches(0.5), cw - Inches(0.3), Inches(0.55),
                 [(b, {"font": UI_FONT, "size": 10, "color": INK})])

    footer(s, total + 1, total + 1)
    return s


def reorder(prs, new_pos_for_appended_indices):
    """Move appended slides into the requested 0-based positions.

    new_pos_for_appended_indices: list of (appended_index, target_position)
    where appended_index counts from start of currently appended slides
    (e.g. 0 = first added, 1 = second added).
    """
    sld_id_lst = prs.slides._sldIdLst
    children = list(sld_id_lst)
    # last len(new_pos_for_appended_indices) entries are the new ones
    n_new = len(new_pos_for_appended_indices)
    new_elems = children[-n_new:]
    base = children[:-n_new]
    # remove all
    for c in children:
        sld_id_lst.remove(c)
    # rebuild
    target = list(base)
    for offset, (added_idx, target_pos) in enumerate(new_pos_for_appended_indices):
        target.insert(target_pos, new_elems[added_idx])
    for c in target:
        sld_id_lst.append(c)


def main():
    prs = Presentation(str(PATH))
    n_before = len(list(prs.slides))
    s_func = add_functionality_slide(prs, n_before + 1)
    s_uniq = add_unique_slide(prs, n_before + 2)

    # Reorder: place these two right after slide 2 (positions 2 and 3 zero-based).
    reorder(prs, [(0, 2), (1, 3)])

    # Re-apply hidden state. New keep-list:
    keep = {1, 2, 3, 4, 9, 12, 13, 17, 19}  # 1-based after reorder:
    # cover=1, vendors=2, NEW func=3, NEW unique=4, status was 8 → now 9 (because 2 inserted at 3,4),
    # but wait: original slides 3-8 shift down by 2. Original 8 → 10. Hmm.
    # Easier: recompute. After reorder, total = n_before + 2. Original slide n>=3 shifts to n+2.
    # Visible originals: 1, 2, 8, 11, 12, 16, 18 → after shift 1, 2, 10, 13, 14, 18, 20 (orig 3..n
    # all shifted by +2 because we inserted 2 new slides at positions 3 and 4)
    # Plus new ones at 3 and 4.
    keep = {1, 2, 3, 4, 10, 13, 14, 18, 20}
    slides = list(prs.slides)
    visible = []
    hidden = []
    for idx, slide in enumerate(slides, 1):
        if idx in keep:
            slide.element.attrib.pop("show", None)
            visible.append(idx)
        else:
            slide.element.set("show", "0")
            hidden.append(idx)

    prs.save(str(PATH))
    print(f"Wrote {PATH}")
    print(f"Visible ({len(visible)}): {visible}")
    print(f"Hidden ({len(hidden)}): {hidden}")


if __name__ == "__main__":
    main()

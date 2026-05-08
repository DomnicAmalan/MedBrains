"""MedBrains chairman + stakeholder deck.

Forest + Copper visual identity. Pitch: do not pay vendor INR 30-70L for
HIS we have already built. Built with python-pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Theme ----------
FOREST       = RGBColor(0x1F, 0x43, 0x32)
FOREST_DEEP  = RGBColor(0x0d, 0x24, 0x17)
FOREST_HOVER = RGBColor(0x15, 0x33, 0x25)
COPPER       = RGBColor(0xB8, 0x92, 0x4A)
INK          = RGBColor(0x0F, 0x14, 0x12)
CANVAS       = RGBColor(0xFF, 0xFF, 0xFF)
PANEL        = RGBColor(0xF7, 0xF8, 0xF6)
RULE         = RGBColor(0xE7, 0xEB, 0xE8)
TINT         = RGBColor(0xE4, 0xED, 0xE9)
DIMMED       = RGBColor(0x55, 0x5C, 0x58)
DANGER       = RGBColor(0xC8, 0x10, 0x2E)
WARN         = RGBColor(0xE6, 0xB4, 0x22)

UI_FONT      = "Inter Tight"
DISP_FONT    = "Fraunces"
MONO_FONT    = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=CANVAS):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_rect(slide, x, y, w, h, fill=FOREST, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text,
    font=UI_FONT, size=14, color=INK, bold=False, italic=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=None,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        runs_data = [(text, {})]
    else:
        runs_data = text
    for i, (txt, opts) in enumerate(runs_data):
        if i == 0:
            run = p.add_run()
        else:
            run = p.add_run()
        run.text = txt
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.italic = opts.get("italic", italic)
        c = opts.get("color", color)
        run.font.color.rgb = c
        if tracking:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def add_eyebrow(slide, x, y, text, color=COPPER, size=10):
    """JetBrains Mono uppercase eyebrow label."""
    return add_text(
        slide, x, y, Inches(8), Inches(0.3),
        text.upper(), font=MONO_FONT, size=size, color=color,
        bold=True, tracking=200,
    )


def add_hairline(slide, x, y, w, color=RULE):
    return add_rect(slide, x, y, w, Emu(9525), fill=color)  # ~0.75pt line


def add_left_accent(slide, x, y, h, color=FOREST, w_pt=3):
    return add_rect(slide, x, y, Pt(w_pt), h, fill=color)


def add_pill(slide, x, y, text, fill=TINT, color=FOREST, w=None, h=None, size=10):
    if w is None:
        w = Inches(1.4)
    if h is None:
        h = Inches(0.32)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.fill.background()
    p.shadow.inherit = False
    tf = p.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = text
    r.font.name = UI_FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return p


def add_footer(slide, page, total, appendix=False):
    add_hairline(slide, Inches(0.5), Inches(7.05), Inches(12.33))
    label = "APPENDIX  ·  Q&A REFERENCE  ·  CONFIDENTIAL" if appendix else \
            "MEDBRAINS HIS  ·  IN-HOUSE BUILD  ·  CONFIDENTIAL"
    add_text(
        slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
        label,
        font=MONO_FONT, size=8,
        color=COPPER if appendix else DIMMED, bold=True, tracking=200,
    )
    add_text(
        slide, Inches(11.5), Inches(7.15), Inches(1.33), Inches(0.3),
        f"{page:02d} / {total:02d}",
        font=MONO_FONT, size=8,
        color=COPPER if appendix else DIMMED, bold=True,
        align=PP_ALIGN.RIGHT, tracking=200,
    )


def add_appendix_badge(slide):
    """Top-right copper pill marking appendix slides."""
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(11.4), Inches(0.45),
                               Inches(1.4), Inches(0.32))
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = COPPER
    p.line.fill.background()
    p.shadow.inherit = False
    tf = p.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = "APPENDIX"
    r.font.name = MONO_FONT; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = CANVAS
    rPr = r._r.get_or_add_rPr(); rPr.set("spc", "200")


def add_brand_mark(slide, x=Inches(0.5), y=Inches(0.45)):
    # small forest square + wordmark
    add_rect(slide, x, y, Inches(0.28), Inches(0.28), fill=FOREST)
    add_text(
        slide, x + Inches(0.4), y - Inches(0.02), Inches(3), Inches(0.35),
        [
            ("Med", {"font": DISP_FONT, "size": 18, "color": INK, "bold": False}),
            ("Brains", {"font": DISP_FONT, "size": 18, "color": FOREST, "bold": False, "italic": True}),
        ],
    )


def big_number(slide, x, y, w, h, value, label, suffix="", color_value=FOREST, color_suffix=COPPER):
    add_text(
        slide, x, y, w, Inches(1.1),
        [
            (str(value), {"font": DISP_FONT, "size": 56, "color": color_value, "bold": False}),
            (suffix, {"font": DISP_FONT, "size": 32, "color": color_suffix, "italic": True}),
        ],
    )
    add_text(
        slide, x, y + Inches(1.05), w, Inches(0.3),
        label.upper(), font=MONO_FONT, size=9, color=DIMMED, bold=True, tracking=200,
    )


def bullets(slide, x, y, w, h, items, size=14, gap=0.32, marker_color=FOREST):
    for i, item in enumerate(items):
        ty = y + Inches(gap * i)
        # forest dot marker
        add_rect(slide, x, ty + Inches(0.1), Pt(6), Pt(6), fill=marker_color)
        add_text(
            slide, x + Inches(0.22), ty, w - Inches(0.22), Inches(gap),
            item, font=UI_FONT, size=size, color=INK,
        )


def section_title(slide, eyebrow, title_runs, subtitle=None):
    add_brand_mark(slide)
    add_eyebrow(slide, Inches(0.5), Inches(1.15), eyebrow)
    add_text(
        slide, Inches(0.5), Inches(1.45), Inches(12.33), Inches(1.3),
        title_runs, anchor=MSO_ANCHOR.TOP,
    )
    if subtitle:
        add_text(
            slide, Inches(0.5), Inches(2.65), Inches(12.33), Inches(0.5),
            subtitle, font=UI_FONT, size=14, color=DIMMED,
        )


# ---------- Slide builders ----------
TOTAL = 24
page = 0
APPENDIX = False


def page_inc():
    global page
    page += 1
    return page


def maybe_appendix(slide):
    if APPENDIX:
        add_appendix_badge(slide)


# 1. Cover
def slide_cover():
    p = page_inc()
    s = add_slide()
    set_bg(s, FOREST)
    # mark
    add_rect(s, Inches(0.6), Inches(0.55), Inches(0.32), Inches(0.32), fill=COPPER)
    add_text(
        s, Inches(1.05), Inches(0.5), Inches(5), Inches(0.5),
        [
            ("Med", {"font": DISP_FONT, "size": 22, "color": CANVAS}),
            ("Brains", {"font": DISP_FONT, "size": 22, "color": COPPER, "italic": True}),
        ],
    )
    add_eyebrow(s, Inches(0.6), Inches(2.4), "Board Briefing  ·  May 2026", color=COPPER)
    add_text(
        s, Inches(0.6), Inches(2.8), Inches(12), Inches(2.4),
        [
            ("The HIS we ", {"font": DISP_FONT, "size": 64, "color": CANVAS}),
            ("already own.", {"font": DISP_FONT, "size": 64, "color": COPPER, "italic": True}),
        ],
        anchor=MSO_ANCHOR.TOP,
    )
    add_text(
        s, Inches(0.6), Inches(4.6), Inches(11), Inches(1.5),
        "Why MedBrains — built in-house — outperforms the ₹30–70 L vendor proposals "
        "on cost, scope, time-to-deploy, and strategic ownership.",
        font=UI_FONT, size=20, color=CANVAS,
    )
    add_hairline(s, Inches(0.6), Inches(6.2), Inches(12), color=COPPER)
    add_text(
        s, Inches(0.6), Inches(6.35), Inches(8), Inches(0.4),
        "FOR THE CHAIRMAN AND HOSPITAL LEADERSHIP",
        font=MONO_FONT, size=9, color=COPPER, bold=True, tracking=240,
    )
    add_text(
        s, Inches(11), Inches(6.35), Inches(2), Inches(0.4),
        "01 / 24",
        font=MONO_FONT, size=9, color=COPPER, bold=True,
        align=PP_ALIGN.RIGHT, tracking=240,
    )


# 2. The decision in front of us
def slide_decision():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "The decision",
        [
            ("Pay a vendor ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("₹30–70 L", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
            (" — or deploy what we have built.", {"font": DISP_FONT, "size": 40, "color": INK}),
        ],
        subtitle="Three vendor proposals on the table. Earliest start: July 2026. "
                 "Meanwhile MedBrains is operational today.",
    )
    # two columns
    col_y = Inches(3.5)
    # Left: vendor
    add_left_accent(s, Inches(0.5), col_y, Inches(2.8), color=DANGER)
    add_text(s, Inches(0.7), col_y, Inches(5.5), Inches(0.4),
             "OPTION A — EXTERNAL VENDOR", font=MONO_FONT, size=10, color=DANGER, bold=True, tracking=200)
    bullets(s, Inches(0.7), col_y + Inches(0.5), Inches(5.7), Inches(2.5), [
        "₹30 L – ₹70 L upfront licence + implementation",
        "AMC 18–22 % per year (₹6–15 L recurring)",
        "Delivery starts July; 6–12 months to go-live",
        "Closed-source binary; vendor owns the IP",
        "Customisation ₹50 K – ₹2 L per change request",
        "Per-hospital re-licence on multi-site expansion",
    ], marker_color=DANGER)
    # Right: us
    add_left_accent(s, Inches(7), col_y, Inches(2.8), color=FOREST)
    add_text(s, Inches(7.2), col_y, Inches(5.5), Inches(0.4),
             "OPTION B — MEDBRAINS (BUILT)", font=MONO_FONT, size=10, color=FOREST, bold=True, tracking=200)
    bullets(s, Inches(7.2), col_y + Inches(0.5), Inches(5.7), Inches(2.5), [
        "₹0 licence cost; runs on hospital infrastructure",
        "Maintenance by internal team only",
        "Operational now — UAT-ready, not a roadmap",
        "Source code + IP held in-house — strategic options open",
        "Customisation same-day, in-house",
        "Unlimited hospitals at zero marginal licence cost",
    ])
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 3. What "₹50 L" actually buys
def slide_vendor_reality():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Vendor reality check",
        [
            ("What the ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("₹50 L quote", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
            (" actually buys.", {"font": DISP_FONT, "size": 40, "color": INK}),
        ],
        subtitle="Standard mid-tier India HIS scope based on Birlamedisoft, Insta, Medixcel, Suvarna comparable proposals.",
    )
    # table
    rows = [
        ("Modules in base licence",      "10–15 (OPD, IPD, Pharmacy, Lab, Billing, HR, MRD basic)", "38 modules ≥80% feature-complete · 62 in active build · 67+ planned"),
        ("Compliance coverage",          "Basic NABH, GST",                                          "NABH+JCI checklists, NDPS, ABDM/NHCX, FHIR R4, GST+TDS, NABL, AERB"),
        ("Patient-safety automation",    "Allergy alerts (basic)",                                   "IPSG, LASA, drug-interaction, dose-validation, consent audit"),
        ("Multi-tenant",                 "Rare; per-instance install",                               "PostgreSQL RLS + 7-layer config native"),
        ("Mobile / TV display apps",     "Add-on, ₹5–15 L extra",                                    "Web + Mobile (RN) + TV (Android TV) in scope"),
        ("Customisation turnaround",     "Weeks; ₹50 K – ₹2 L per change",                           "Same-day, in-house"),
        ("Source code access",           "Closed",                                                   "Held in-house — open-source / internal / commercial path open"),
    ]
    x = Inches(0.5); y = Inches(3.45); w = Inches(12.33)
    add_rect(s, x, y, w, Inches(0.4), fill=FOREST)
    add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(4.2), Inches(0.3),
             "DIMENSION", font=MONO_FONT, size=9, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(4.4), y + Inches(0.08), Inches(4.2), Inches(0.3),
             "VENDOR", font=MONO_FONT, size=9, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(8.6), y + Inches(0.08), Inches(4.2), Inches(0.3),
             "MEDBRAINS", font=MONO_FONT, size=9, color=COPPER, bold=True, tracking=200)
    rh = Inches(0.42)
    for i, (dim, ven, us) in enumerate(rows):
        ry = y + Inches(0.4) + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.2), ry + Inches(0.08), Inches(4.2), rh, dim,
                 font=UI_FONT, size=11, color=INK, bold=True)
        add_text(s, x + Inches(4.4), ry + Inches(0.08), Inches(4.1), rh, ven,
                 font=UI_FONT, size=10, color=DIMMED)
        add_text(s, x + Inches(8.6), ry + Inches(0.08), Inches(4.2), rh, us,
                 font=UI_FONT, size=10, color=FOREST, bold=True)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 4. The hidden math — vendor 5-yr TCO
def slide_tco():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Five-year total cost",
        [
            ("Vendor TCO is not ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("the quote.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
        subtitle="A ₹50 L proposal compounds into ₹1.25 Cr+ over five years — before any second hospital.",
    )
    # bar comparison
    base_x = Inches(0.7); base_y = Inches(6.4); max_h = Inches(2.7)
    bar_w = Inches(1.0); gap = Inches(0.18)
    items = [
        ("Year 0 impl.",    50, COPPER),
        ("AMC × 5 yr",      50, COPPER),
        ("Custom dev",      15, COPPER),
        ("2nd hospital",    35, COPPER),
        ("Vendor 5-yr TCO", 150, DANGER),
        ("MedBrains 5-yr",  0,  FOREST),
    ]
    max_val = 150
    for i, (label, val, color) in enumerate(items):
        bx = base_x + (bar_w + gap) * i
        h = Emu(int(max_h * (val / max_val))) if val > 0 else Inches(0.05)
        add_rect(s, bx, base_y - h, bar_w, h, fill=color)
        add_text(s, bx, base_y - h - Inches(0.4), bar_w, Inches(0.3),
                 f"₹{val} L" if val > 0 else "₹0",
                 font=DISP_FONT, size=18, color=color, align=PP_ALIGN.CENTER)
        add_text(s, bx, base_y + Inches(0.05), bar_w, Inches(0.5),
                 label, font=UI_FONT, size=10, color=INK, align=PP_ALIGN.CENTER)
    # baseline
    add_hairline(s, Inches(0.5), base_y, Inches(8.5), color=RULE)
    # right summary
    add_left_accent(s, Inches(9.5), Inches(3.5), Inches(2.8), color=COPPER)
    add_text(s, Inches(9.7), Inches(3.5), Inches(3.3), Inches(0.4),
             "WHAT THE BAR DOES NOT SHOW", font=MONO_FONT, size=9, color=COPPER, bold=True, tracking=200)
    bullets(s, Inches(9.7), Inches(3.95), Inches(3.4), Inches(2.5), [
        "Productivity loss during 6–12 mo rollout",
        "Re-training cost on vendor exit",
        "Data export / migration on switch",
        "Lock-in premium on AMC renegotiation",
        "Per-site licence on hospital expansion",
    ], size=11, gap=0.36, marker_color=COPPER)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 5. What we built — at a glance
def slide_built_overview():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "What we have built",
        [
            ("A working hospital platform — ", {"font": DISP_FONT, "size": 38, "color": INK}),
            ("today.", {"font": DISP_FONT, "size": 38, "color": COPPER, "italic": True}),
        ],
        subtitle="Numbers are measured against the working repository, not promised.",
    )
    y = Inches(3.4)
    cols = [
        ("547,960", "lines of code", " "),
        ("77", "database migrations", " "),
        ("125", "backend modules", " "),
        ("138", "frontend pages", " "),
    ]
    cw = Inches(3.0); gap = Inches(0.15); x = Inches(0.5)
    for i, (val, label, suf) in enumerate(cols):
        bx = x + (cw + gap) * i
        add_rect(s, bx, y, cw, Inches(2.0), fill=PANEL)
        add_left_accent(s, bx, y, Inches(2.0), color=FOREST)
        big_number(s, bx + Inches(0.25), y + Inches(0.35), cw - Inches(0.3), Inches(1.4),
                   val, label, suffix=suf)
    # second row — capabilities
    y2 = Inches(5.6)
    caps = [
        ("48.5%", "feature catalogue covered (1,596 of 3,291)"),
        ("38", "modules ≥80% feature-complete"),
        ("11", "regulatory frameworks integrated"),
        ("3", "platforms — Web, Mobile, TV"),
    ]
    for i, (val, label) in enumerate(caps):
        bx = x + (cw + gap) * i
        add_text(s, bx, y2, cw, Inches(0.6),
                 [(val, {"font": DISP_FONT, "size": 28, "color": FOREST})],
                 anchor=MSO_ANCHOR.TOP)
        add_text(s, bx, y2 + Inches(0.55), cw, Inches(0.5),
                 label, font=UI_FONT, size=11, color=DIMMED)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 6. Module coverage map
def slide_module_map():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Module coverage",
        [
            ("Thirty-eight modules ", {"font": DISP_FONT, "size": 38, "color": INK}),
            ("live.", {"font": DISP_FONT, "size": 38, "color": COPPER, "italic": True}),
        ],
        subtitle="Vendor packages typically ship 10–15. We are already past three vendor scopes.",
    )
    groups = [
        ("Clinical core",
         ["OPD", "IPD", "ICU", "Emergency", "OT", "Doctor Dashboard", "CPOE Order Sets", "Care View", "Bedside Portal"]),
        ("Diagnostics",
         ["Lab (Phase 1-3)", "Radiology", "Blood Bank", "CSSD", "BME"]),
        ("Pharmacy & Supply",
         ["Pharmacy (Phase 2)", "NDPS Register", "Indent / Inventory (Phase 2)", "Procurement", "Stores & Transfers"]),
        ("Patient & Front Office",
         ["Patient Mgmt", "Appointments", "Front Office & Queue", "Camp Management", "Communications / DLT"]),
        ("Finance",
         ["Billing (Phase 3)", "GST / TDS / GSTR", "Bank Reconciliation", "Journal & GL", "Insurance / TPA / NHCX"]),
        ("Workforce & Quality",
         ["HR & Staff", "Quality (NABH)", "Infection Control", "Consent Mgmt", "MRD", "Audit"]),
        ("Operations & Compliance",
         ["Facilities (MGPS, Fire, Water, Energy)", "Housekeeping", "Diet & Kitchen", "Regulatory", "Occ. Health", "ABDM / FHIR"]),
        ("Platform",
         ["Multi-Hospital", "Print Engine", "Dashboard Builder", "Screen Builder", "TV Displays", "Onboarding & Setup"]),
    ]
    cols = 2; rows = 4
    cw = Inches(6.1); rh = Inches(0.85); x0 = Inches(0.5); y0 = Inches(3.3); gx = Inches(0.13); gy = Inches(0.1)
    for i, (title, mods) in enumerate(groups):
        c = i % cols; r = i // cols
        bx = x0 + (cw + gx) * c
        by = y0 + (rh + gy) * r
        add_rect(s, bx, by, cw, rh, fill=PANEL)
        add_left_accent(s, bx, by, rh, color=FOREST)
        add_text(s, bx + Inches(0.2), by + Inches(0.08), cw, Inches(0.3),
                 title.upper(), font=MONO_FONT, size=9, color=FOREST, bold=True, tracking=200)
        add_text(s, bx + Inches(0.2), by + Inches(0.4), cw - Inches(0.25), Inches(0.45),
                 "  ·  ".join(mods), font=UI_FONT, size=10, color=INK)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 7. Feature comparison matrix
def slide_feature_matrix():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Feature comparison",
        [
            ("Vendor scope vs ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("MedBrains scope.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
        subtitle="Tick = available out of the box. Cross = quoted as add-on or out of scope.",
    )
    rows = [
        ("OPD, IPD, Pharmacy, Lab, Billing core",          True,  True),
        ("ICU bedside, OT scheduling, Emergency triage",   False, True),
        ("Blood Bank, CSSD, Diet, BME",                    False, True),
        ("NDPS register · Schedule H/H1/X enforcement",    False, True),
        ("ABDM / NHCX integration · Health-ID linking",    False, True),
        ("FHIR R4 export · DICOM · HL7 v2",                False, True),
        ("GST, TDS, GSTR returns, bank reconciliation",    False, True),
        ("Multi-tenant · 7-layer hierarchical config",     False, True),
        ("Mobile app (RN) and TV displays in scope",       False, True),
        ("NABH + JCI checklists, audit, evidence engine",  False, True),
        ("Source code held in-house · strategic options preserved", False, True),
    ]
    x = Inches(0.5); y = Inches(3.4); w = Inches(12.33)
    rh = Inches(0.32)
    add_rect(s, x, y, w, Inches(0.36), fill=FOREST)
    add_text(s, x + Inches(0.2), y + Inches(0.05), Inches(8.5), Inches(0.3),
             "CAPABILITY", font=MONO_FONT, size=9, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(8.9), y + Inches(0.05), Inches(1.7), Inches(0.3),
             "VENDOR", font=MONO_FONT, size=9, color=CANVAS, bold=True, tracking=200, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(10.6), y + Inches(0.05), Inches(1.7), Inches(0.3),
             "MEDBRAINS", font=MONO_FONT, size=9, color=COPPER, bold=True, tracking=200, align=PP_ALIGN.CENTER)
    for i, (cap, v, m) in enumerate(rows):
        ry = y + Inches(0.36) + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.2), ry + Inches(0.05), Inches(8.5), rh, cap,
                 font=UI_FONT, size=11, color=INK)
        add_text(s, x + Inches(8.9), ry + Inches(0.05), Inches(1.7), rh,
                 "✓" if v else "✗",
                 font=UI_FONT, size=14,
                 color=FOREST if v else DANGER,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(10.6), ry + Inches(0.05), Inches(1.7), rh,
                 "✓" if m else "✗",
                 font=UI_FONT, size=14,
                 color=FOREST if m else DANGER,
                 bold=True, align=PP_ALIGN.CENTER)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 8. Build value
def slide_build_value():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Build value realised",
        [
            ("What this asset is ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("worth.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
        subtitle="Five independent valuation methods. All point upward of the highest vendor quote.",
    )
    methods = [
        ("Per-feature method",
         "₹59 L",
         "1,360 Done × ₹4 K + 236 Partial × ₹2 K, India dev rate."),
        ("Per-module method",
         "₹1.9 Cr",
         "38 modules × ₹5 L mid-tier outsourced build cost."),
        ("LOC × ₹40 (conservative)",
         "₹2.19 Cr",
         "Production-grade code at the lower end of India rates."),
        ("LOC × ₹100 (skilled)",
         "₹5.48 Cr",
         "Skilled full-stack rate including tests and migrations."),
        ("Vendor TCO replaced (5 yr)",
         "₹1.25 Cr",
         "₹50 L impl + ₹50 L AMC + ₹15 L customisation + ₹10 L per extra site."),
    ]
    x = Inches(0.5); y = Inches(3.4); rh = Inches(0.7); w = Inches(12.33)
    for i, (name, val, expl) in enumerate(methods):
        ry = y + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_left_accent(s, x, ry, rh, color=COPPER if i == 4 else FOREST)
        add_text(s, x + Inches(0.25), ry + Inches(0.12), Inches(4.5), Inches(0.3),
                 name, font=UI_FONT, size=12, color=INK, bold=True)
        add_text(s, x + Inches(0.25), ry + Inches(0.4), Inches(4.5), Inches(0.3),
                 expl, font=UI_FONT, size=10, color=DIMMED)
        add_text(s, x + Inches(8.5), ry + Inches(0.15), Inches(3.5), Inches(0.5),
                 val, font=DISP_FONT, size=28,
                 color=COPPER if i == 4 else FOREST,
                 align=PP_ALIGN.RIGHT)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 9. Clinical depth
def slide_clinical():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Clinical depth",
        [
            ("Care covered end-to-end — ", {"font": DISP_FONT, "size": 38, "color": INK}),
            ("OPD to discharge.", {"font": DISP_FONT, "size": 38, "color": COPPER, "italic": True}),
        ],
    )
    cards = [
        ("OPD",
         "Queue, doctor dashboard, CPOE order sets, e-prescription, follow-ups, repeat-Rx."),
        ("IPD",
         "Admission, bed-state, structured nursing notes, MAR, vitals, restraint, discharge workflow."),
        ("ICU",
         "Bedside vitals, device pairing, handoff, family communication, structured rounds."),
        ("Emergency",
         "Triage, MLC documentation, code blue/red/pink workflow, disposition, time-stamped audit."),
        ("OT",
         "Scheduling, anaesthesia, surgical consent, OT analytics — surgeon caseload, complications."),
        ("Lab",
         "Phlebotomy queue, batch QC (Levey-Jennings + Westgard), critical alerts, NABL records, B2B."),
        ("Radiology",
         "Order, report, PACS-ready, AERB compliance hooks, PCPNDT lockdown."),
        ("Pharmacy",
         "OTC, in-patient, discharge, NDPS dual register, batch FEFO, formulary, AWaRe stewardship."),
    ]
    cw = Inches(3.0); rh = Inches(1.55); gx = Inches(0.13); gy = Inches(0.13); x0 = Inches(0.5); y0 = Inches(3.3)
    for i, (h, b) in enumerate(cards):
        c = i % 4; r = i // 4
        bx = x0 + (cw + gx) * c
        by = y0 + (rh + gy) * r
        add_rect(s, bx, by, cw, rh, fill=PANEL)
        add_left_accent(s, bx, by, rh, color=FOREST)
        add_text(s, bx + Inches(0.2), by + Inches(0.15), cw, Inches(0.4),
                 h, font=DISP_FONT, size=18, color=FOREST)
        add_text(s, bx + Inches(0.2), by + Inches(0.55), cw - Inches(0.3), Inches(0.95),
                 b, font=UI_FONT, size=10, color=INK)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 10. Compliance built-in
def slide_compliance():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Compliance built-in",
        [
            ("Regulation is ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("not a roadmap item.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
        subtitle="Indian healthcare law and accreditation are enforced in code, not bolted on.",
    )
    items = [
        ("NDPS Act 1985",            "Dual-lock register, dispensing audit, daily reconciliation."),
        ("D&C Act / Schedules H, H1, X", "Schedule-aware dispensing checks; Rx duplicate enforcement."),
        ("PCPNDT Act",               "Sex-determination block, audit trail on every imaging order."),
        ("NABH + JCI",               "34 department checklists, 700+ criteria, evidence engine."),
        ("ABDM / NHCX",              "Health-ID linking, claims exchange, callback log."),
        ("HL7 FHIR R4 + DICOM + HL7 v2", "Interoperability ready; not a future task."),
        ("GST + TDS + GSTR + HSN",   "Returns generation, TDS deposit, certificate, ERP export."),
        ("PNDT, MTP, MHCA, BMW 2016","Domain rules embedded in module business logic."),
    ]
    x = Inches(0.5); y = Inches(3.4); rh = Inches(0.45); w = Inches(12.33)
    for i, (h, b) in enumerate(items):
        ry = y + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.25), ry + Inches(0.1), Inches(3.8), rh,
                 h, font=UI_FONT, size=12, color=INK, bold=True)
        add_text(s, x + Inches(4.2), ry + Inches(0.1), Inches(8), rh,
                 b, font=UI_FONT, size=11, color=DIMMED)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 11. Patient safety
def slide_safety():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Patient safety",
        [
            ("Six IPSGs, ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("automated.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    goals = [
        ("IPSG 1", "Patient identification (UHID + 2nd identifier on every action)"),
        ("IPSG 2", "Effective communication — read-back, SBAR, structured handoff"),
        ("IPSG 3", "High-alert medication safety — LASA flags, dose ranges by age/weight/renal"),
        ("IPSG 4", "Surgical safety — site marking, time-out, sign-out checklists"),
        ("IPSG 5", "Infection prevention — hand-hygiene, HAI surveillance, antibiogram"),
        ("IPSG 6", "Falls prevention — risk score, reassessment, alert escalation"),
    ]
    x = Inches(0.5); y = Inches(3.4); rh = Inches(0.55); w = Inches(12.33)
    for i, (g, b) in enumerate(goals):
        ry = y + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_left_accent(s, x, ry, rh, color=COPPER)
        add_text(s, x + Inches(0.25), ry + Inches(0.13), Inches(1.6), rh,
                 g, font=DISP_FONT, size=18, color=FOREST)
        add_text(s, x + Inches(2.0), ry + Inches(0.18), Inches(10.2), rh,
                 b, font=UI_FONT, size=12, color=INK)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 12. Interoperability
def slide_interop():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Interoperability",
        [
            ("Speaks every ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("standard.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
        subtitle="No vendor portal lock-in. Data exits cleanly; partners plug in cleanly.",
    )
    cards = [
        ("ABDM / NHCX",     "Health-ID, consent manager, claims exchange callback log."),
        ("HL7 FHIR R4",     "Patient, Encounter, Observation, MedicationRequest, Claim."),
        ("DICOM",           "Imaging instance routing; PACS-ready storage hooks."),
        ("HL7 v2",          "Lab analyser interfaces, ADT messaging."),
        ("LOINC",           "Lab test catalogue keyed for cross-lab interop."),
        ("ICD-10 / 11",     "Diagnoses with WHO coding."),
        ("WHO ATC + INN",   "Drug catalogue with international nonproprietary names."),
        ("RxNorm-ready",    "Cross-mapping for global drug master."),
    ]
    cw = Inches(3.0); rh = Inches(1.4); gx = Inches(0.13); gy = Inches(0.13); x0 = Inches(0.5); y0 = Inches(3.4)
    for i, (h, b) in enumerate(cards):
        c = i % 4; r = i // 4
        bx = x0 + (cw + gx) * c
        by = y0 + (rh + gy) * r
        add_rect(s, bx, by, cw, rh, fill=CANVAS, line=RULE)
        add_text(s, bx + Inches(0.2), by + Inches(0.15), cw, Inches(0.4),
                 h, font=DISP_FONT, size=18, color=FOREST)
        add_text(s, bx + Inches(0.2), by + Inches(0.55), cw - Inches(0.3), Inches(0.8),
                 b, font=UI_FONT, size=10, color=INK)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 13. Tech stack advantage
def slide_tech():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Architecture",
        [
            ("Modern stack vs ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("legacy vendor stack.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    rows = [
        ("Backend language",  "PHP / .NET / legacy Java (most India HIS)", "Rust 2024 — memory safe, no nulls, compile-time SQL"),
        ("Web framework",     "Server-rendered, jQuery / WebForms",        "Axum 0.8 + Tower, async, typed routes"),
        ("Frontend",          "Bootstrap 3, jQuery",                       "React 18 + Mantine v7 + TanStack Query v5"),
        ("Database",          "MySQL / MSSQL, single-tenant",              "PostgreSQL 16 with RLS + YottaDB for hierarchical config"),
        ("Multi-tenant",      "Per-instance install, no isolation",         "Native RLS, transaction-scoped tenant context"),
        ("Auth",              "Session cookies, MD5 hashes",               "Ed25519 JWT, Argon2id, 743-permission RBAC across 55 modules"),
        ("CI / Tests",        "Manual QA",                                 "Compile-time SQL checks, Playwright E2E, contract checks"),
        ("Cross-platform",    "Web only or paid mobile add-on",            "Web + React Native CLI mobile + Android TV — one codebase tree"),
    ]
    x = Inches(0.5); y = Inches(3.3); w = Inches(12.33); rh = Inches(0.42)
    add_rect(s, x, y, w, Inches(0.36), fill=FOREST)
    add_text(s, x + Inches(0.2), y + Inches(0.05), Inches(3.8), Inches(0.3),
             "LAYER", font=MONO_FONT, size=9, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(4.1), y + Inches(0.05), Inches(4.2), Inches(0.3),
             "TYPICAL VENDOR HIS", font=MONO_FONT, size=9, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(8.4), y + Inches(0.05), Inches(4.2), Inches(0.3),
             "MEDBRAINS", font=MONO_FONT, size=9, color=COPPER, bold=True, tracking=200)
    for i, (l, v, m) in enumerate(rows):
        ry = y + Inches(0.36) + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.2), ry + Inches(0.1), Inches(3.8), rh, l,
                 font=UI_FONT, size=11, color=INK, bold=True)
        add_text(s, x + Inches(4.1), ry + Inches(0.1), Inches(4.2), rh, v,
                 font=UI_FONT, size=10, color=DIMMED)
        add_text(s, x + Inches(8.4), ry + Inches(0.1), Inches(4.2), rh, m,
                 font=UI_FONT, size=10, color=FOREST, bold=True)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 14. Security
def slide_security():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Security",
        [
            ("Defence in depth — ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("not a checkbox.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    items = [
        ("Database row-level security",  "PostgreSQL RLS enforced at every transaction. No accidental cross-tenant leak."),
        ("Authentication",                "Ed25519 JWT, Argon2id password hashing — both modern, audited primitives."),
        ("Authorisation",                 "743 typed permissions across 55 modules; per-user override (extra / denied)."),
        ("Audit trail",                   "Append-only audit table, partitioned, signed where regulatory requires."),
        ("Break-glass + scope expiry",    "Emergency access controlled, time-limited, logged."),
        ("Data residency",                "Self-hosted; runs on hospital infrastructure or chosen Indian region."),
        ("Encryption",                    "TLS in transit, at-rest via PostgreSQL TDE/disk encryption per deployment."),
        ("Backup + DR",                   "Idempotent migrations, point-in-time recovery, ARM starter + multi-region path."),
    ]
    x = Inches(0.5); y = Inches(3.4); rh = Inches(0.45); w = Inches(12.33)
    for i, (h, b) in enumerate(items):
        ry = y + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.25), ry + Inches(0.1), Inches(3.8), rh,
                 h, font=UI_FONT, size=12, color=INK, bold=True)
        add_text(s, x + Inches(4.2), ry + Inches(0.1), Inches(8), rh,
                 b, font=UI_FONT, size=11, color=DIMMED)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 15. Three platforms
def slide_platforms():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Three platforms",
        [
            ("Web, mobile, TV — ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("one codebase tree.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    cards = [
        ("Web", "Mantine v7 + SCSS",
         "Doctor, nurse, admin, billing, pharmacy, lab — full ops UI. 138 pages."),
        ("Mobile", "React Native CLI + Paper v5",
         "Doctor rounds, nurse charting, patient bedside portal, ambulance triage. New Architecture, offline-first WatermelonDB."),
        ("TV", "React Native Android TV",
         "OPD queue boards, IPD bed status dashboards, OT scheduler, kitchen display, digital signage."),
    ]
    cw = Inches(4.05); rh = Inches(3.3); gx = Inches(0.15); x0 = Inches(0.5); y0 = Inches(3.4)
    for i, (t, sub, body) in enumerate(cards):
        bx = x0 + (cw + gx) * i
        add_rect(s, bx, y0, cw, rh, fill=PANEL)
        add_left_accent(s, bx, y0, rh, color=FOREST)
        add_text(s, bx + Inches(0.25), y0 + Inches(0.2), cw, Inches(0.6),
                 t, font=DISP_FONT, size=32, color=FOREST)
        add_text(s, bx + Inches(0.25), y0 + Inches(0.95), cw, Inches(0.3),
                 sub.upper(), font=MONO_FONT, size=9, color=COPPER, bold=True, tracking=200)
        add_text(s, bx + Inches(0.25), y0 + Inches(1.4), cw - Inches(0.4), rh - Inches(1.5),
                 body, font=UI_FONT, size=12, color=INK)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 16. Customisation speed
def slide_speed():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Speed of change",
        [
            ("A request becomes a feature in ", {"font": DISP_FONT, "size": 36, "color": INK}),
            ("hours, not weeks.", {"font": DISP_FONT, "size": 36, "color": COPPER, "italic": True}),
        ],
    )
    # left vendor timeline, right ours
    add_left_accent(s, Inches(0.5), Inches(3.4), Inches(3.4), color=DANGER)
    add_text(s, Inches(0.7), Inches(3.4), Inches(5.5), Inches(0.4),
             "VENDOR CHANGE REQUEST", font=MONO_FONT, size=10, color=DANGER, bold=True, tracking=200)
    bullets(s, Inches(0.7), Inches(3.9), Inches(5.5), Inches(3), [
        "Day 0 — Raise change request",
        "Day 3–7 — Vendor scoping + quote (₹50 K – ₹2 L)",
        "Day 7–14 — PO approval cycle",
        "Day 14–35 — Vendor development + QA",
        "Day 35–45 — UAT + deploy on next release window",
        "Outcome: 4–6 weeks per non-trivial change",
    ], size=12, gap=0.36, marker_color=DANGER)
    add_left_accent(s, Inches(7), Inches(3.4), Inches(3.4), color=FOREST)
    add_text(s, Inches(7.2), Inches(3.4), Inches(5.5), Inches(0.4),
             "MEDBRAINS CHANGE REQUEST", font=MONO_FONT, size=10, color=FOREST, bold=True, tracking=200)
    bullets(s, Inches(7.2), Inches(3.9), Inches(5.5), Inches(3), [
        "Hour 0 — Department raises request",
        "Hour 1 — In-house team triages",
        "Hour 2–8 — Build + tests (Playwright + Rust clippy)",
        "Hour 8 — Staging deploy + UAT",
        "Hour 24 — Production rollout",
        "Outcome: same-day to next-day per change",
    ], size=12, gap=0.36)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 17. Vendor risks
def slide_vendor_risks():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Vendor risks",
        [
            ("Five risks the vendor proposal does ", {"font": DISP_FONT, "size": 36, "color": INK}),
            ("not price.", {"font": DISP_FONT, "size": 36, "color": COPPER, "italic": True}),
        ],
    )
    risks = [
        ("Lock-in",        "Closed source, vendor-specific schema. Switching costs near-equal to a new build."),
        ("AMC creep",      "Renegotiation power sits with the vendor. 18–22 % per year compounds."),
        ("Scope erosion",  "Mid-tier proposals exclude ICU, OT, NDPS, ABDM, FHIR — billed later or not at all."),
        ("Delivery slip",  "July start ≠ July go-live. Industry mean rollout: 9 months. Cash flows out before value lands."),
        ("Compliance gap", "Vendor MVPs rarely include NDPS dual register, NABH evidence engine, GSTR generation. Hospital absorbs liability."),
    ]
    x = Inches(0.5); y = Inches(3.4); rh = Inches(0.62); w = Inches(12.33)
    for i, (h, b) in enumerate(risks):
        ry = y + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_left_accent(s, x, ry, rh, color=DANGER)
        add_text(s, x + Inches(0.25), ry + Inches(0.13), Inches(2.5), rh,
                 h, font=UI_FONT, size=14, color=DANGER, bold=True)
        add_text(s, x + Inches(3), ry + Inches(0.18), Inches(9.2), rh,
                 b, font=UI_FONT, size=11, color=INK)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 18. Timeline comparison
def slide_timeline():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Timeline",
        [
            ("Vendor starts in July. We are ", {"font": DISP_FONT, "size": 36, "color": INK}),
            ("operational now.", {"font": DISP_FONT, "size": 36, "color": COPPER, "italic": True}),
        ],
    )
    # horizontal timeline
    x0 = Inches(0.7); y0 = Inches(4.1); w = Inches(11.9)
    add_hairline(s, x0, y0, w, color=RULE)
    months = ["May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec", "Jan'27", "Feb'27", "Mar'27"]
    seg = w / (len(months) - 1)
    for i, m in enumerate(months):
        cx = x0 + seg * i
        add_rect(s, cx - Pt(2), y0 - Pt(3), Pt(4), Pt(8), fill=DIMMED)
        add_text(s, cx - Inches(0.4), y0 + Inches(0.15), Inches(0.8), Inches(0.3),
                 m, font=MONO_FONT, size=9, color=DIMMED, align=PP_ALIGN.CENTER)
    # vendor band
    vy = y0 - Inches(1.3)
    vstart = x0 + seg * 2  # July
    vend = x0 + seg * 10
    add_rect(s, vstart, vy, vend - vstart, Inches(0.45), fill=DANGER)
    add_text(s, vstart, vy - Inches(0.32), Inches(6), Inches(0.3),
             "VENDOR — KICKOFF · DEV · UAT · GO-LIVE",
             font=MONO_FONT, size=9, color=DANGER, bold=True, tracking=200)
    add_text(s, vstart + Inches(0.15), vy + Inches(0.1), Inches(6), Inches(0.3),
             "₹50 L paid; value delivered ~Mar 2027",
             font=UI_FONT, size=11, color=CANVAS, bold=True)
    # ours band
    my = y0 - Inches(2.4)
    add_rect(s, x0, my, w, Inches(0.45), fill=FOREST)
    add_text(s, x0, my - Inches(0.32), Inches(6), Inches(0.3),
             "MEDBRAINS — IN PRODUCTION TODAY · CONTINUOUS DELIVERY",
             font=MONO_FONT, size=9, color=FOREST, bold=True, tracking=200)
    add_text(s, x0 + Inches(0.15), my + Inches(0.1), Inches(8), Inches(0.3),
             "UAT-ready now; new modules ship every fortnight",
             font=UI_FONT, size=11, color=CANVAS, bold=True)
    # legend below
    add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
             "Every month of vendor wait is a month our patients still run on paper.",
             font=DISP_FONT, size=18, color=INK, italic=True)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 19. Cost comparison final
def slide_cost_final():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Five-year cost — final",
        [
            ("₹1.25 Cr saved. ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("Per hospital.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    rows = [
        ("Implementation / licence",            "₹50,00,000",   "₹0"),
        ("AMC × 5 years (10 L / yr)",           "₹50,00,000",   "₹0"),
        ("Customisation budget (5 yr)",         "₹15,00,000",   "Internal team only"),
        ("Second hospital licence",             "₹35,00,000",   "₹0 — unlimited"),
        ("Mobile / TV add-on",                  "₹10,00,000",   "Included"),
        ("Compliance gap remediation",          "₹8,00,000",    "Built-in"),
    ]
    x = Inches(0.5); y = Inches(3.3); w = Inches(12.33); rh = Inches(0.46)
    add_rect(s, x, y, w, Inches(0.4), fill=FOREST)
    add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(5), Inches(0.3),
             "LINE ITEM", font=MONO_FONT, size=10, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(5.5), y + Inches(0.08), Inches(3.5), Inches(0.3),
             "VENDOR", font=MONO_FONT, size=10, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(9), y + Inches(0.08), Inches(3.5), Inches(0.3),
             "MEDBRAINS", font=MONO_FONT, size=10, color=COPPER, bold=True, tracking=200)
    for i, (l, v, m) in enumerate(rows):
        ry = y + Inches(0.4) + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.2), ry + Inches(0.13), Inches(5), rh, l,
                 font=UI_FONT, size=12, color=INK, bold=True)
        add_text(s, x + Inches(5.5), ry + Inches(0.13), Inches(3.5), rh, v,
                 font=UI_FONT, size=12, color=DANGER, bold=True)
        add_text(s, x + Inches(9), ry + Inches(0.13), Inches(3.5), rh, m,
                 font=UI_FONT, size=12, color=FOREST, bold=True)
    # totals row
    ty = y + Inches(0.4) + rh * len(rows)
    add_rect(s, x, ty, w, Inches(0.55), fill=FOREST)
    add_text(s, x + Inches(0.2), ty + Inches(0.15), Inches(5), Inches(0.4),
             "FIVE-YEAR TCO (PER HOSPITAL)", font=MONO_FONT, size=11, color=CANVAS, bold=True, tracking=200)
    add_text(s, x + Inches(5.5), ty + Inches(0.1), Inches(3.5), Inches(0.5),
             "₹1.68 Cr", font=DISP_FONT, size=22, color=COPPER)
    add_text(s, x + Inches(9), ty + Inches(0.1), Inches(3.5), Inches(0.5),
             "₹0 licence", font=DISP_FONT, size=22, color=CANVAS)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 20. Strategic options — management decides
def slide_strategic():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Strategic options",
        [
            ("Three paths — ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("management decides.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
        subtitle="Adopting MedBrains today does not lock in any path. The platform is built; "
                 "how it is positioned commercially is a separate decision the board can take later.",
    )
    paths = [
        ("Path A",
         "Internal use only",
         [
            "Run MedBrains across our own hospital(s).",
            "No external distribution, no licensing.",
            "Lowest disclosure, full control of roadmap.",
            "Cost saving = vendor TCO replaced.",
         ]),
        ("Path B",
         "Open-source release",
         [
            "Publish under a permissive or copyleft licence.",
            "Build community, attract contributors, set the standard.",
            "Reduces long-term maintenance burden.",
            "Brand value as a healthcare-tech leader.",
         ]),
        ("Path C",
         "Commercial product",
         [
            "Licence MedBrains to peer hospitals / chains.",
            "Recurring revenue stream — SaaS or per-bed.",
            "Spin-out vehicle possible; investor optionality.",
            "Competes directly with the vendors quoting us today.",
         ]),
    ]
    x = Inches(0.5); y = Inches(3.5); cw = Inches(4.05); gx = Inches(0.15)
    accent = [FOREST, COPPER, FOREST_DEEP]
    for i, (tag, name, items) in enumerate(paths):
        bx = x + (cw + gx) * i
        add_rect(s, bx, y, cw, Inches(3.3), fill=PANEL)
        add_left_accent(s, bx, y, Inches(3.3), color=accent[i], w_pt=4)
        add_text(s, bx + Inches(0.25), y + Inches(0.2), cw, Inches(0.4),
                 tag.upper(), font=MONO_FONT, size=10, color=accent[i], bold=True, tracking=200)
        add_text(s, bx + Inches(0.25), y + Inches(0.55), cw, Inches(0.5),
                 name, font=DISP_FONT, size=22, color=INK)
        bullets(s, bx + Inches(0.25), y + Inches(1.25), cw - Inches(0.4), Inches(2.0),
                items, size=11, gap=0.42, marker_color=accent[i])
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4),
             "Today's decision: adopt internally. Path B and Path C remain open and need no commitment now.",
             font=UI_FONT, size=11, color=DIMMED, italic=True)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 21. Roadmap
def slide_roadmap():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Roadmap",
        [
            ("Already planned, ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("not invented on the call.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    phases = [
        ("Now — Q2 2026",
         ["UAT in two pilot departments", "Doctor + nurse training rolling", "Pharmacy NDPS register cut-over",
          "ABDM Health-ID enrolment", "Mobile app rollout to senior consultants"]),
        ("Q3 2026",
         ["YottaDB hot path enabled (real-time bed state, sequences)", "Specialty modules — Psychiatry, Onco workflows",
          "TV displays in OPD + IPD", "Multi-hospital tenancy onboarding", "Quality evidence engine for NABH cycle"]),
        ("Q4 2026 — H1 2027",
         ["AI-assisted documentation (notes, discharge summary)", "Predictive analytics — readmission, sepsis early warning",
          "PACS integration cutover", "ABDM HCX claims at scale", "Medical College / academic module rollout"]),
    ]
    x = Inches(0.5); y = Inches(3.3); cw = Inches(4.05); gx = Inches(0.15)
    for i, (ph, items) in enumerate(phases):
        bx = x + (cw + gx) * i
        add_rect(s, bx, y, cw, Inches(3.5), fill=PANEL)
        add_left_accent(s, bx, y, Inches(3.5), color=FOREST if i == 0 else (COPPER if i == 2 else FOREST_HOVER))
        add_text(s, bx + Inches(0.25), y + Inches(0.2), cw, Inches(0.4),
                 ph.upper(), font=MONO_FONT, size=10, color=FOREST, bold=True, tracking=200)
        bullets(s, bx + Inches(0.25), y + Inches(0.7), cw - Inches(0.4), Inches(2.7),
                items, size=11, gap=0.42)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 22. Risk + answers
def slide_qa():
    p = page_inc()
    s = add_slide(); set_bg(s)
    section_title(
        s, "Anticipated questions",
        [
            ("What the board will ", {"font": DISP_FONT, "size": 40, "color": INK}),
            ("ask — answered.", {"font": DISP_FONT, "size": 40, "color": COPPER, "italic": True}),
        ],
    )
    qa = [
        ("Bus-factor — what if the in-house team leaves?",
         "Source code, tests, migrations and RFCs are version-controlled and documented. A successor team onboards faster than a vendor handover."),
        ("How is this different from vendor X who also said multi-tenant?",
         "Show me the row-level policy in their schema. We can show ours, today, in the migration history."),
        ("Will it pass NABH?",
         "We mapped 700+ NABH/JCI criteria into evidence engine and module checklists. NABH-ready before vendor would finish OPD."),
        ("What about insurance / TPA?",
         "ABDM/NHCX live, dual-insurance, reimbursement docs, TPA recon already in billing module."),
        ("Hardware integration — analysers, printers, barcode?",
         "Architecture supports it. Lab analyser HL7 v2 + DICOM hooks built. Pending items are integration projects, not blocked rebuilds."),
        ("Risk if we say no to the vendor today?",
         "Zero. Vendor will requote in 30 days. We will be three modules ahead."),
    ]
    x = Inches(0.5); y = Inches(3.3); rh = Inches(0.55); w = Inches(12.33)
    for i, (q, a) in enumerate(qa):
        ry = y + rh * i
        bg = PANEL if i % 2 == 0 else CANVAS
        add_rect(s, x, ry, w, rh, fill=bg)
        add_text(s, x + Inches(0.25), ry + Inches(0.08), Inches(5.5), rh,
                 "Q  " + q, font=UI_FONT, size=11, color=INK, bold=True)
        add_text(s, x + Inches(6), ry + Inches(0.08), Inches(6.3), rh,
                 "A  " + a, font=UI_FONT, size=10, color=FOREST)
    maybe_appendix(s); add_footer(s, p, TOTAL, appendix=APPENDIX)


# 23. Recommendation
# 6. Closing — recommendation + the ask, combined
def slide_close():
    p = page_inc()
    s = add_slide(); set_bg(s, FOREST)
    add_brand_mark(s)
    add_eyebrow(s, Inches(0.5), Inches(1.15), "Recommendation & ask", color=COPPER)
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(2),
             [
                 ("Decline the vendor proposal. ", {"font": DISP_FONT, "size": 38, "color": CANVAS}),
                 ("Adopt MedBrains.", {"font": DISP_FONT, "size": 38, "color": COPPER, "italic": True}),
             ], anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.5), Inches(2.6), Inches(12), Inches(0.5),
             "Three decisions. One meeting.",
             font=UI_FONT, size=18, color=CANVAS)
    asks = [
        ("01", "Endorsement",
         "Board endorses MedBrains as the HIS in use. Vendor proposals declined or repositioned to integration-only."),
        ("02", "Budget",
         "Approve internal-team operating budget for FY26 — delivery, infra, training. Capex avoided ≈ ₹50 L vs vendor."),
        ("03", "Sponsor",
         "Name an executive sponsor and pilot department leads. UAT begins within 14 days of approval."),
    ]
    x = Inches(0.5); y = Inches(3.6); cw = Inches(4.05); gx = Inches(0.15)
    for i, (n, h, b) in enumerate(asks):
        bx = x + (cw + gx) * i
        add_rect(s, bx, y, cw, Inches(2.7), fill=CANVAS)
        add_left_accent(s, bx, y, Inches(2.7), color=COPPER, w_pt=4)
        add_text(s, bx + Inches(0.25), y + Inches(0.2), cw, Inches(0.5),
                 n, font=DISP_FONT, size=32, color=COPPER)
        add_text(s, bx + Inches(0.25), y + Inches(0.85), cw, Inches(0.4),
                 h, font=DISP_FONT, size=20, color=FOREST)
        add_text(s, bx + Inches(0.25), y + Inches(1.3), cw - Inches(0.4), Inches(1.3),
                 b, font=UI_FONT, size=11, color=INK)
    add_hairline(s, Inches(0.5), Inches(7.05), Inches(12.33), color=COPPER)
    add_text(s, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             "MEDBRAINS HIS  ·  IN-HOUSE BUILD  ·  CONFIDENTIAL",
             font=MONO_FONT, size=8, color=COPPER, bold=True, tracking=200)
    add_text(s, Inches(11.5), Inches(7.15), Inches(1.33), Inches(0.3),
             f"{p:02d} / {TOTAL:02d}",
             font=MONO_FONT, size=8, color=COPPER, bold=True,
             align=PP_ALIGN.RIGHT, tracking=200)


# 7. Thank you / Q&A divider
def slide_thanks():
    p = page_inc()
    s = add_slide(); set_bg(s)
    add_brand_mark(s)
    add_eyebrow(s, Inches(0.5), Inches(2.6), "End of main deck")
    add_text(s, Inches(0.5), Inches(2.95), Inches(12.33), Inches(2.5),
             [
                 ("Thank you. ", {"font": DISP_FONT, "size": 80, "color": INK}),
                 ("Questions?", {"font": DISP_FONT, "size": 80, "color": COPPER, "italic": True}),
             ])
    add_hairline(s, Inches(0.5), Inches(5.5), Inches(12.33), color=COPPER)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(0.5),
             "The slides that follow are reference material for the discussion.",
             font=UI_FONT, size=16, color=DIMMED, italic=True)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.5),
             "Vendor reality · TCO breakdown · Module map · Feature matrix · Build value · "
             "Clinical · Compliance · Patient safety · Interop · Architecture · Security · "
             "Platforms · Speed · Risks · Timeline · Roadmap · Q&A.",
             font=UI_FONT, size=11, color=DIMMED)
    add_hairline(s, Inches(0.5), Inches(7.05), Inches(12.33))
    add_text(s, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             "MEDBRAINS HIS  ·  IN-HOUSE BUILD  ·  CONFIDENTIAL",
             font=MONO_FONT, size=8, color=DIMMED, bold=True, tracking=200)
    add_text(s, Inches(11.5), Inches(7.15), Inches(1.33), Inches(0.3),
             f"{p:02d} / {TOTAL:02d}",
             font=MONO_FONT, size=8, color=DIMMED, bold=True,
             align=PP_ALIGN.RIGHT, tracking=200)


def main():
    global APPENDIX
    # === MAIN — 6 slides ===
    slide_cover()
    slide_decision()
    slide_built_overview()
    slide_cost_final()
    slide_strategic()
    slide_close()
    # === Thank-you divider ===
    slide_thanks()
    # === APPENDIX — 17 reference slides ===
    APPENDIX = True
    slide_vendor_reality()
    slide_tco()
    slide_module_map()
    slide_feature_matrix()
    slide_build_value()
    slide_clinical()
    slide_compliance()
    slide_safety()
    slide_interop()
    slide_tech()
    slide_security()
    slide_platforms()
    slide_speed()
    slide_vendor_risks()
    slide_timeline()
    slide_roadmap()
    slide_qa()
    out = "/Users/apple/Projects/MedBrains/presentations/MedBrains-Chairman-Deck.pptx"
    prs.save(out)
    print(f"Wrote {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

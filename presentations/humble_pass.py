"""Soften deck tone — drop know-it-all phrasing, keep facts.

Applies a humility pass to the seven visible slides of Alagappa_Pitch_v2.pptx.
Hidden slides are left untouched.
"""
from pathlib import Path

from pptx import Presentation

PATH = Path("/Users/apple/Projects/MedBrains/presentations/Alagappa_Pitch_v2.pptx")


def replace_run_starting_with(slide, prefix: str, new_text: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.startswith(prefix):
                    run.text = new_text
                    return True
    return False


def replace_run_exact(slide, old: str, new: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == old:
                    run.text = new
                    return True
    return False


def replace_run_contains(slide, needle: str, new: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if needle in run.text:
                    run.text = new
                    return True
    return False


def collapse_paragraph(slide, full_old: str, new: str) -> bool:
    """When a paragraph's runs combine to ``full_old``, replace into one run."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            joined = "".join(r.text for r in para.runs)
            if joined == full_old and para.runs:
                para.runs[0].text = new
                for r in para.runs[1:]:
                    r.text = ""
                return True
    return False


def soften(prs: Presentation) -> list[tuple[int, str, bool]]:
    log: list[tuple[int, str, bool]] = []
    slides = list(prs.slides)

    # --- Slide 1 ---
    s = slides[0]
    ok = replace_run_exact(
        s,
        "Pay ₹14L–₹70L upfront to a vendor — or spend a fraction this year and own it.",
        "Pay ₹14L–₹70L upfront to a vendor — or spend a fraction this year and keep full control of the build.",
    )
    log.append((1, "cover line", ok))

    # --- Slide 2 ---
    s = slides[1]
    ok = replace_run_exact(
        s,
        " Each one a different gamble.",
        " Each with different trade-offs.",
    )
    log.append((2, "different gamble → trade-offs", ok))

    # --- Slide 8 ---
    s = slides[7]
    ok1 = replace_run_exact(s, "WHY ME · WHY NOW", "WHERE THE BUILD STANDS")
    log.append((8, "eyebrow WHY ME → WHERE THE BUILD STANDS", ok1))
    ok2 = replace_run_exact(s, "I've already ", "Phase 1 scope is ")
    log.append((8, "I've already → Phase 1 scope is", ok2))
    ok3 = replace_run_exact(s, "built", "substantially")
    log.append((8, "built → substantially", ok3))
    ok4 = replace_run_exact(s, " 80% of it.", " covered today.")
    log.append((8, "80% of it. → covered today.", ok4))
    ok5 = replace_run_exact(s, "STACK MASTERY", "TECHNICAL STACK")
    log.append((8, "STACK MASTERY → TECHNICAL STACK", ok5))

    # --- Slide 11 ---
    s = slides[10]
    ok = collapse_paragraph(
        s,
        "Even at similar 3-yr cost, in-house gives: source ownership · zero AMC after Y3 · resale revenue · ₹0 marginal branch expansion.",
        "At similar 3-yr cost, the in-house path keeps full source-code control, no AMC after Y3, no per-branch licence, and an optional resale path the board can choose later.",
    )
    if not ok:
        ok = replace_run_contains(
            s,
            "source ownership",
            "At similar 3-yr cost, the in-house path keeps full source-code control, no AMC after Y3, no per-branch licence, and an optional resale path the board can choose later.",
        )
    log.append((11, "footer line", ok))

    # --- Slide 12 ---
    s = slides[11]
    ok = collapse_paragraph(
        s,
        "Pre-launch saving alone: ~₹42 L. And we own the platform from day one — vendor lock-in starts at ₹0.",
        "Pre-launch saving alone: ~₹42 L. And no vendor lock-in clock starts ticking — the platform stays under our roof.",
    )
    if not ok:
        ok = replace_run_contains(
            s,
            "we own the platform",
            "Pre-launch saving alone: ~₹42 L. And no vendor lock-in clock starts ticking — the platform stays under our roof.",
        )
    log.append((12, "footer line", ok))

    # --- Slide 16 ---
    s = slides[15]
    ok1 = replace_run_exact(
        s,
        "Begin Y2 sales prep · first sister-hospital deploy",
        "Optional — explore first sister-hospital deployment",
    )
    log.append((16, "Oct'26 line", ok1))
    ok2 = replace_run_exact(
        s,
        "Y2 reseller revenue begins · 2 partner hospitals signed",
        "Optional — reseller path begins, if the board chooses",
    )
    log.append((16, "Q2'27 line", ok2))

    return log


def main() -> None:
    prs = Presentation(str(PATH))
    log = soften(prs)
    prs.save(str(PATH))
    print(f"Wrote {PATH}")
    for sl, what, ok in log:
        flag = "✓" if ok else "✗ NOT FOUND"
        print(f"  S{sl:>2}  {flag}  {what}")


if __name__ == "__main__":
    main()

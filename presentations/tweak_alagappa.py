"""Apply small, factual tweaks to Alagappa_Pitch.pptx.

Two kinds of edit:
1. Update stale counts to match the audited repo state (May 2026).
2. Soften the closing ask — drop the "Approve the budget / team / timeline"
   framing in favour of a humbler "where we stand / what's needed / suggested
   next step" wording.

All visual layout, fonts, colours, ordering and other content are preserved.
Output goes alongside the original under presentations/.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

SRC = Path("/Users/apple/Downloads/Alagappa_Pitch.pptx")
OUT = Path("/Users/apple/Projects/MedBrains/presentations/Alagappa_Pitch_v2.pptx")

# Run-level replacements applied across the whole deck. Order matters; longer
# strings are tried first so that "1,360 / 3,291" is not partially matched by a
# later rule.
GLOBAL_REPLACEMENTS: list[tuple[str, str]] = [
    ("998 / 2,566", "1,360 / 3,291"),
    ("998/2,566", "1,360/3,291"),
    ("250k+", "547k+"),
    ("250K+", "547K+"),
]


def apply_global_replacements(prs: Presentation) -> int:
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in GLOBAL_REPLACEMENTS:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            n += 1
    return n


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


def update_run_in_text_frame(tf, old: str, new: str) -> bool:
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text.strip() == old.strip():
                run.text = new
                return True
    return False


def find_text_frame_containing(slide, needle: str):
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            return shape.text_frame
    return None


def patch_slide_8_modules_metric(slide) -> bool:
    """Slide 8 — '44' tile (next to 'modules production-ready') → '38'.

    The two are in separate text boxes on the slide, so we identify the right
    '44' by checking the slide has a 'modules production-ready' shape and
    update the lone-'44' shape.
    """
    has_label = any(
        shape.has_text_frame and "modules production-ready" in shape.text_frame.text
        for shape in slide.shapes
    )
    if not has_label:
        return False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        if len(runs) == 1 and runs[0].text.strip() == "44":
            runs[0].text = "38"
            return True
    return False


def patch_slide_9_subtitle(slide) -> bool:
    """Slide 9 — '44 production modules.' → '38 production modules.'"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.startswith("44 production modules"):
                    run.text = run.text.replace(
                        "44 production modules", "38 production modules"
                    )
                    return True
    return False


def _replace_run_text(slide, old: str, new: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == old:
                    run.text = new
                    return True
    return False


def patch_slide_16_subtitle(slide) -> bool:
    """Slide 16 — 'Approve in May → ' → 'May start → ' (split-run aware)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.startswith("Approve in May"):
                    run.text = "May start → "
                    return True
    return False


def patch_slide_18_ask(slide) -> int:
    """Slide 18 — rewrite the three 'Approve …' cards to a humbler frame.

    The slide's runs match exactly, so we replace text in-place keeping the
    fonts/sizes intact.
    """
    edits: list[tuple[str, str]] = [
        ("THE ASK", "WHERE WE STAND"),
        ("Three approvals. This week.",
         "Status, what's needed, and a suggested next step."),
        # card 1
        ("Approve the budget", "What's already done"),
        ("Modest 3-year envelope. Hardware + operations only. Versus ₹50–70 L upfront + AMC forever to a vendor.",
         "18 months of build. 547K+ lines of code. 1,360 features live in dev across 38 production-ready modules. NHCX, ABDM, FHIR R4, NDPS, GST/TDS — code complete."),
        # card 2
        ("Approve the team", "What it would take to finish"),
        ("Hire 2 senior engineers. I lead. We start in June. The team becomes Alagappa's long-term tech capability.",
         "Two senior engineers (Rust + React) joining the existing in-house lead. Hardware engineer already on staff. The same team can maintain the platform after go-live."),
        # card 3
        ("Approve the timeline", "Suggested next step"),
        ("Pilot in July. Full go-live by hospital opening day. First reseller signed by Q2 2027.",
         "A two-week pilot in July would let us validate the OPD and Pharmacy paths side-by-side with current process before the September opening. Decision rests with the board."),
        # footer line
        ("Save ₹40 L+ pre-launch.  Earn ₹15 cr over 5 yrs.  Own the platform.",
         "Pre-launch saving ≈ ₹40 L vs the average vendor quote. Resale upside is optional, not assumed."),
    ]
    n = 0
    for old, new in edits:
        if _replace_run_text(slide, old, new):
            n += 1
        else:
            # try splitting into pieces in case the run was broken up by formatting
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    full = "".join(r.text for r in para.runs)
                    if full == old and para.runs:
                        para.runs[0].text = new
                        for r in para.runs[1:]:
                            r.text = ""
                        n += 1
                        break
    return n


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"Source not found: {SRC}")
    prs = Presentation(str(SRC))
    slides = list(prs.slides)
    if len(slides) < 18:
        raise SystemExit(f"Expected 18 slides, got {len(slides)}")

    n_global = apply_global_replacements(prs)
    s8_done = patch_slide_8_modules_metric(slides[7])
    s9_done = patch_slide_9_subtitle(slides[8])
    s16_done = patch_slide_16_subtitle(slides[15])
    s18_n = patch_slide_18_ask(slides[17])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    print(f"Global replacements: {n_global}")
    print(f"Slide 8 modules metric: {'updated' if s8_done else 'NOT found'}")
    print(f"Slide 9 subtitle: {'updated' if s9_done else 'NOT found'}")
    print(f"Slide 16 subtitle: {'updated' if s16_done else 'NOT found'}")
    print(f"Slide 18 ask edits: {s18_n} / 9")


if __name__ == "__main__":
    main()
